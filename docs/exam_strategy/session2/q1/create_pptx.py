#!/usr/bin/env python3
"""
Generate PowerPoint for Q1 exam strategy.
- Slide 1: Title
- Slide 2: Verbatim question & options
- Slides 3-7: One per option (diagram + explanation)
- Slide 8: Answer key summary
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "q1_exam_strategy.pptx")

# ── Colors ──
AWS_DARK    = RGBColor(0x23, 0x2F, 0x3E)
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CORRECT_GRN = RGBColor(0x2E, 0x7D, 0x32)
INCORRECT_R = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY    = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT    = RGBColor(0x21, 0x21, 0x21)
SUBTLE_BG   = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, font_size=13,
                       color=DARK_TEXT, font_name="Calibri", line_spacing=None,
                       bold_flags=None):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        if bold_flags and i < len(bold_flags):
            p.font.bold = bold_flags[i]
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), AWS_ORANGE)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
             "AWS Data Migration Strategy",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
             "Hybrid Environment  |  5 PB Historical  |  3 TB/month Ongoing  |  Amazon S3 + Redshift",
             font_size=20, color=AWS_ORANGE, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), AWS_ORANGE)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "Exam Strategy  —  Question 1  —  Select TWO correct answers",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.43), SLIDE_W, Inches(0.07), AWS_ORANGE)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Verbatim Question & Options
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

# Header
add_shape(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75), AWS_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.65),
             "QUESTION",
             font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

# Scenario text — verbatim
scenario = (
    "Your customer is an enterprise sized business in the telecommunications industry. "
    "They are looking to setup a hybrid environment and intend to utilize AWS Redshift "
    "for querying a dataset comprising 5 PB of historical cell phone usage data. They "
    "have an ongoing data accumulation rate of 3 TB per month. The requirement is for "
    "the new incoming monthly data to be accessible for querying with less than a "
    "two-day delay. You are responsible for determining the most efficient method to "
    "transfer this data into Amazon S3."
)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.3),
             scenario, font_size=13, color=DARK_TEXT, line_spacing=20)

# "Which combination..." line
add_text_box(slide, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.4),
             "Which combination of steps will provide an appropriate solution? (Select TWO)",
             font_size=14, color=AWS_DARK, bold=True)

# Options — verbatim, each in its own card
options_raw = [
    ("A.", "Leverage AWS Direct Connect for both the initial bulk transfer of historical data and the ongoing monthly data transfers, ensuring consistent network performance."),
    ("B.", "Implement AWS DataSync for continuous synchronization of the 3 TB monthly data to Amazon S3, maintaining the required freshness of less than two days."),
    ("C.", "Configure Amazon S3 Transfer Acceleration for the ongoing 3 TB monthly data uploads to minimize transfer times and meet the delay requirements."),
    ("D.", "Use AWS Snowball for the initial migration of the 5 PB of historical data to Amazon S3, ensuring a fast and secure data transfer method for large datasets."),
    ("E.", "Directly upload the ongoing 3 TB of data each month via the Internet using the AWS CLI or SDKs, relying on your existing internet bandwidth."),
]

card_y = Inches(3.3)
card_h = Inches(0.72)
card_gap = Inches(0.08)

for i, (letter, text) in enumerate(options_raw):
    y = card_y + i * (card_h + card_gap)

    add_shape(slide, Inches(0.8), y, Inches(11.7), card_h, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))

    # Letter badge — uniform gray for all options
    add_shape(slide, Inches(1.0), y + Inches(0.12), Inches(0.48), Inches(0.48), MED_GRAY)
    add_text_box(slide, Inches(1.0), y + Inches(0.12), Inches(0.48), Inches(0.48),
                 letter.replace(".", ""), font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Option text
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(10.5), card_h - Inches(0.15),
                 text, font_size=12, color=DARK_TEXT)

# Timer GIF — bottom right of question slide
timer_gif = os.path.join(os.path.dirname(BASE), "countdown_2min.gif")
if os.path.exists(timer_gif):
    slide.shapes.add_picture(timer_gif, Inches(10.5), Inches(6.5), Inches(2.5), Inches(0.85))


# ══════════════════════════════════════════════════════════════
# SLIDES 3-7: Option detail slides (diagram + explanation)
# ══════════════════════════════════════════════════════════════
options_detail = [
    {
        "file": "option_a_direct_connect.png",
        "letter": "A",
        "title": "AWS Direct Connect",
        "verdict": "INCORRECT",
        "is_correct": False,
        # Each bullet: (text, is_key_reason)
        "bullets": [
            ("Provides a dedicated 1-400 Gbps private network connection to AWS", False),
            ("Primarily designed for low-latency hybrid access, not bulk data migration", True),
            ("Snowball and DataSync are purpose-built for the 5 PB bulk and 3 TB/month recurring requirements", True),
            ("Requires physical cross-connect setup at a colocation facility — adds weeks of lead time", False),
        ],
    },
    {
        "file": "option_b_datasync.png",
        "letter": "B",
        "title": "AWS DataSync",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("Transfers data from on-premises (NFS/SMB/HDFS) to Amazon S3 or EFS", False),
            ("Supports continuous synchronization with built-in scheduling — ideal for 3 TB/month", True),
            ("Purpose-built protocol delivers up to 10 Gbps per task with TLS encryption", False),
            ("Meets the < 2-day freshness requirement through automated, recurring transfers", True),
            ("Includes in-transit and at-rest integrity validation", False),
        ],
    },
    {
        "file": "option_c_transfer_acceleration.png",
        "letter": "C",
        "title": "S3 Transfer Acceleration",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Speeds up uploads by routing data through 50+ CloudFront edge locations", False),
            ("No built-in scheduling or incremental sync — must be managed manually", True),
            ("DataSync is better suited for continuous synchronization and the < 2-day SLA", True),
            ("Best for: geographically distant, ad-hoc large uploads — not recurring sync", False),
        ],
    },
    {
        "file": "option_d_snowball.png",
        "letter": "D",
        "title": "AWS Snowball",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("Physical storage devices (~210 TB each) shipped to your data center", False),
            ("Optimized for large, one-time migrations — 5 PB is a perfect use case", True),
            ("Physical transfer bypasses all network bottlenecks — faster than any link at PB scale", True),
            ("256-bit encryption with tamper-resistant enclosure and TPM chip", False),
            ("AWS handles shipping, S3 import, and device sanitization", False),
        ],
    },
    {
        "file": "option_e_direct_upload.png",
        "letter": "E",
        "title": "Direct Internet Upload (CLI/SDK)",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Relies entirely on existing internet bandwidth — often limited for enterprises", True),
            ("No transfer optimization, compression, or automatic retry on failure", False),
            ("At 100 Mbps, 3 TB takes ~2.8 days — exceeds the < 2-day SLA", True),
            ("DataSync provides scheduling, integrity checks, and optimized throughput natively", False),
        ],
    },
]


for idx, opt in enumerate(options_detail):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    accent_color = CORRECT_GRN if opt["is_correct"] else INCORRECT_R

    # Top accent
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent_color)
    # Bottom accent
    add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

    # Header bar
    add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.75), AWS_DARK)

    # Letter badge in header
    add_shape(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55), accent_color)
    add_text_box(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55),
                 opt["letter"], font_size=24, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(0.38), Inches(8), Inches(0.6),
                 opt["title"], font_size=24, color=WHITE, bold=True)

    # Verdict badge
    add_shape(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5), accent_color)
    add_text_box(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5),
                 opt["verdict"], font_size=16, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # ── Option text card (below header, above diagram) ──
    option_text = options_raw[idx][1]
    opt_card_top = Inches(1.16)
    opt_card_h = Inches(0.65)

    add_shape(slide, Inches(0.53), opt_card_top, Inches(12.27), opt_card_h,
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))

    add_shape(slide, Inches(0.73), opt_card_top + Inches(0.1),
              Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(0.73), opt_card_top + Inches(0.1),
                 Inches(0.45), Inches(0.45),
                 opt["letter"], font_size=16, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.43), opt_card_top + Inches(0.1),
                 Inches(11.17), Inches(0.30),
                 option_text, font_size=12, color=DARK_TEXT)

    # ── Diagram: full-width, shifted down for option card ──
    img_path = os.path.join(BASE, opt["file"])
    img = Image.open(img_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    diagram_area_left = Inches(0.5)
    diagram_area_top = Inches(1.95)
    diagram_area_w = Inches(12.3)
    diagram_area_h = Inches(3.25)

    if diagram_area_w / diagram_area_h > aspect:
        final_h = diagram_area_h
        final_w = int(diagram_area_h * aspect)
    else:
        final_w = diagram_area_w
        final_h = int(diagram_area_w / aspect)

    img_left = diagram_area_left + (diagram_area_w - final_w) // 2
    img_top = diagram_area_top + (diagram_area_h - final_h) // 2

    slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    # ── Explanation panel: horizontal bar at bottom ──
    panel_top = Inches(5.36)
    panel_h = Inches(1.94)

    add_shape(slide, Inches(0.5), panel_top, Inches(12.3), panel_h,
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))

    # Accent stripe
    add_rect(slide, Inches(0.5), panel_top + Inches(0.19), Inches(0.06), Inches(1.65), accent_color)

    # Verdict label
    verdict_label = "Correct" if opt["is_correct"] else "Incorrect"
    add_text_box(slide, Inches(0.85), panel_top + Inches(0.15),
                 Inches(1.5), Inches(0.35),
                 verdict_label, font_size=16, color=accent_color, bold=True)

    # Separator
    add_text_box(slide, Inches(2.2), panel_top + Inches(0.15),
                 Inches(0.3), Inches(0.35),
                 "|", font_size=16, color=RGBColor(0xE0, 0xE0, 0xE0), bold=True)

    add_text_box(slide, Inches(2.45), panel_top + Inches(0.15),
                 Inches(1.5), Inches(0.35),
                 "Explanation", font_size=14, color=AWS_DARK, bold=True)

    # Thin line
    add_rect(slide, Inches(0.85), panel_top + Inches(0.55),
             Inches(11.7), Inches(0.015), RGBColor(0xE0, 0xE0, 0xE0))

    # Bullet-point explanation
    bullets = opt["bullets"]
    txBox = slide.shapes.add_textbox(
        Inches(0.85), panel_top + Inches(0.64), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for bi, (bullet_text, is_key) in enumerate(bullets):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.line_spacing = Pt(18)
        p.font.name = "Calibri"

        if is_key:
            run_m = p.add_run()
            run_m.text = "\u25B6  "
            run_m.font.size = Pt(9)
            run_m.font.color.rgb = accent_color
            run_m.font.bold = True
            run_m.font.name = "Calibri"

            run_t = p.add_run()
            run_t.text = bullet_text
            run_t.font.size = Pt(11)
            run_t.font.color.rgb = accent_color
            run_t.font.bold = True
            run_t.font.name = "Calibri"
        else:
            run_m = p.add_run()
            run_m.text = "\u2022  "
            run_m.font.size = Pt(11)
            run_m.font.color.rgb = MED_GRAY
            run_m.font.name = "Calibri"

            run_t = p.add_run()
            run_t.text = bullet_text
            run_t.font.size = Pt(11)
            run_t.font.color.rgb = DARK_TEXT
            run_t.font.bold = False
            run_t.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Answer Key Summary
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "ANSWER KEY",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.2), Inches(1.3), Inches(2.9), Inches(0.04), AWS_ORANGE)

answers = [
    ("A", "AWS Direct Connect",           "INCORRECT", INCORRECT_R,
     "Dedicated private link — not optimized for bulk PB migration or scheduled recurring sync."),
    ("B", "AWS DataSync",                  "CORRECT",   CORRECT_GRN,
     "Scheduled sync with purpose-built protocol (10 Gbps), integrity checks — ideal for 3 TB/month."),
    ("C", "S3 Transfer Acceleration",      "INCORRECT", INCORRECT_R,
     "Edge-routed uploads — no built-in scheduling or incremental sync capability."),
    ("D", "AWS Snowball",                  "CORRECT",   CORRECT_GRN,
     "Physical devices (~210 TB each) — fastest and most practical for 5 PB one-time migration."),
    ("E", "Direct Upload (CLI/SDK)",       "INCORRECT", INCORRECT_R,
     "Public internet — unreliable, no optimization, 100 Mbps would take ~2.8 days exceeding SLA."),
]

for i, (letter, name, verdict, color, reason) in enumerate(answers):
    y = Inches(1.8) + i * Inches(1.0)

    row_bg = RGBColor(0x2C, 0x3A, 0x4A) if i % 2 == 0 else RGBColor(0x34, 0x44, 0x55)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), row_bg)

    add_shape(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5), color)
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 letter, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(3.2), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)

    add_shape(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45), color)
    add_text_box(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45),
                 verdict, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.65),
                 reason, font_size=11, color=LIGHT_GRAY)

# Bottom answer callout
add_rect(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45), AWS_ORANGE)
add_text_box(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45),
             "Correct Answers:  B (DataSync for 3 TB/month)  +  D (Snowball for 5 PB historical)",
             font_size=15, color=AWS_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_ORANGE)


# ── Save ──
prs.save(OUT)
print(f"Saved: {OUT}")
