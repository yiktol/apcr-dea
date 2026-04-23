#!/usr/bin/env python3
"""Generate PowerPoint for Q5 exam strategy."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "q5_exam_strategy.pptx")

AWS_DARK    = RGBColor(0x23, 0x2F, 0x3E)
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CORRECT_GRN = RGBColor(0x2E, 0x7D, 0x32)
INCORRECT_R = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY    = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT    = RGBColor(0x21, 0x21, 0x21)
SUBTLE_BG   = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color: s.line.color.rgb = border_color; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_rect(slide, left, top, width, height, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color; s.line.fill.background()
    return s

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font_name; p.alignment = alignment
    if line_spacing: p.line_spacing = Pt(line_spacing)
    return txBox


# ═══ SLIDE 1: Title ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), AWS_ORANGE)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
             "S3 Storage Class Optimization",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
             "S3 Standard  |  Standard-IA  |  Glacier Instant Retrieval  |  Glacier Deep Archive",
             font_size=20, color=AWS_ORANGE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), AWS_ORANGE)
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "Exam Strategy  —  Question 5  —  Most cost-effective combination",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.43), SLIDE_W, Inches(0.07), AWS_ORANGE)


# ═══ SLIDE 2: Verbatim Question ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

add_shape(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75), AWS_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.65),
             "QUESTION", font_size=24, color=WHITE, bold=True)

scenario = (
    "A data engineer works at a video streaming service that stores large amounts of "
    "video content on Amazon S3. They have a mix of frequently accessed content that "
    "needs to be readily available as well as a large amount of older content that is "
    "accessed once per quarter. The data engineer is tasked with optimizing storage "
    "costs while maintaining the required performance for their workload"
)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1),
             scenario, font_size=13, color=DARK_TEXT, line_spacing=20)

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.4),
             "Which combination of Amazon S3 storage classes will meet the requirements in the most cost effective way?",
             font_size=14, color=AWS_DARK, bold=True)

options_raw = [
    ("A.", "S3 Standard for frequently accessed content and S3 Glacier Instant Retrieval for infrequently accessed content."),
    ("B.", "S3 Standard for frequently accessed content and S3 Infrequent-Access for infrequently accessed content."),
    ("C.", "S3 Glacier Instant Retrieval for frequently accessed content and S3 Glacier Deep Archive for infrequently accessed archive content"),
    ("D.", "S3 Standard-IA for frequently accessed content, S3 Glacier Instant Retrieval for infrequently accessed archive content."),
]

card_y = Inches(3.2)
card_h = Inches(0.72)
card_gap = Inches(0.08)

for i, (letter, text) in enumerate(options_raw):
    y = card_y + i * (card_h + card_gap)
    add_shape(slide, Inches(0.8), y, Inches(11.7), card_h, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45),
                 letter.replace(".", ""), font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(10.5), card_h - Inches(0.15),
                 text, font_size=12, color=DARK_TEXT)

# Timer GIF
timer_gif = os.path.join(os.path.dirname(BASE), "countdown_2min.gif")
if os.path.exists(timer_gif):
    slide.shapes.add_picture(timer_gif, Inches(10.5), Inches(6.5), Inches(2.5), Inches(0.85))


# ═══ SLIDES 3-6: Option details ═══
options_detail = [
    {
        "file": "option_a_standard_glacier_ir.png",
        "letter": "A",
        "title": "S3 Standard + Glacier Instant Retrieval",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("S3 Standard for frequently accessed content — no retrieval fee, millisecond access", False),
            ("Glacier Instant Retrieval saves 68% vs Standard-IA for data accessed once per quarter", True),
            ("Glacier Instant Retrieval still provides millisecond access — no latency trade-off", True),
            ("Higher per-GB retrieval cost offset by much lower storage cost at quarterly access frequency", False),
            ("AWS explicitly recommends this class for data accessed approximately once per quarter", False),
        ],
    },
    {
        "file": "option_b_standard_ia.png",
        "letter": "B",
        "title": "S3 Standard + Standard-IA",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("S3 Standard for hot content is correct", False),
            ("Standard-IA works for infrequent access but costs 68% more than Glacier Instant Retrieval at quarterly access", True),
            ("Standard-IA is better suited for monthly access patterns, not once-per-quarter", True),
            ("Tempting answer — functional but not the MOST cost-effective option", False),
        ],
    },
    {
        "file": "option_c_glacier_both.png",
        "letter": "C",
        "title": "Glacier Instant + Glacier Deep Archive",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Glacier Instant Retrieval is NOT designed for frequently accessed content — high retrieval costs", True),
            ("Glacier Deep Archive has 12-48 hour retrieval — far too slow for quarterly video streaming", True),
            ("Frequently accessed content requires S3 Standard for zero retrieval fees and high throughput", False),
            ("Deep Archive is for data retained 7-10 years and accessed less than once per year", False),
        ],
    },
    {
        "file": "option_d_ia_glacier_ir.png",
        "letter": "D",
        "title": "Standard-IA + Glacier Instant Retrieval",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Standard-IA is NOT suitable for frequently accessed content — $0.01/GB retrieval fee on every access", True),
            ("Frequent access on Standard-IA accumulates high retrieval costs quickly", True),
            ("Glacier Instant Retrieval for quarterly content is correct, but the hot tier is wrong", False),
            ("S3 Standard is the only appropriate class for frequently accessed data — no retrieval fee", False),
        ],
    },
]

for idx, opt in enumerate(options_detail):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    accent_color = CORRECT_GRN if opt["is_correct"] else INCORRECT_R

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent_color)
    add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.75), AWS_DARK)
    add_shape(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55), accent_color)
    add_text_box(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55),
                 opt["letter"], font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(0.38), Inches(8), Inches(0.6),
                 opt["title"], font_size=24, color=WHITE, bold=True)
    add_shape(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5), accent_color)
    add_text_box(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5),
                 opt["verdict"], font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Option text card
    option_text = options_raw[idx][1]
    oc_top = Inches(1.16)
    add_shape(slide, Inches(0.53), oc_top, Inches(12.27), Inches(0.65),
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, Inches(0.73), oc_top + Inches(0.1), Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(0.73), oc_top + Inches(0.1), Inches(0.45), Inches(0.45),
                 opt["letter"], font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.43), oc_top + Inches(0.1), Inches(11.17), Inches(0.30),
                 option_text, font_size=12, color=DARK_TEXT)

    # Diagram
    img_path = os.path.join(BASE, opt["file"])
    img = Image.open(img_path)
    aspect = img.size[0] / img.size[1]
    da_l, da_t, da_w, da_h = Inches(0.5), Inches(1.95), Inches(12.3), Inches(3.25)
    if da_w / da_h > aspect:
        fh = da_h; fw = int(da_h * aspect)
    else:
        fw = da_w; fh = int(da_w / aspect)
    slide.shapes.add_picture(img_path, da_l + (da_w - fw)//2, da_t + (da_h - fh)//2, fw, fh)

    # Explanation panel
    pt = Inches(5.36); ph = Inches(1.94)
    add_shape(slide, Inches(0.5), pt, Inches(12.3), ph, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, Inches(0.5), pt + Inches(0.19), Inches(0.06), Inches(1.65), accent_color)

    vl = "Correct" if opt["is_correct"] else "Incorrect"
    add_text_box(slide, Inches(0.85), pt + Inches(0.15), Inches(1.5), Inches(0.35),
                 vl, font_size=16, color=accent_color, bold=True)
    add_text_box(slide, Inches(2.2), pt + Inches(0.15), Inches(0.3), Inches(0.35),
                 "|", font_size=16, color=RGBColor(0xE0, 0xE0, 0xE0), bold=True)
    add_text_box(slide, Inches(2.45), pt + Inches(0.15), Inches(1.5), Inches(0.35),
                 "Explanation", font_size=14, color=AWS_DARK, bold=True)
    add_rect(slide, Inches(0.85), pt + Inches(0.55), Inches(11.7), Inches(0.015),
             RGBColor(0xE0, 0xE0, 0xE0))

    txBox = slide.shapes.add_textbox(Inches(0.85), pt + Inches(0.64), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame; tf.word_wrap = True
    for bi, (bt, is_key) in enumerate(opt["bullets"]):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.line_spacing = Pt(18); p.font.name = "Calibri"
        if is_key:
            rm = p.add_run(); rm.text = "\u25B6  "; rm.font.size = Pt(9)
            rm.font.color.rgb = accent_color; rm.font.bold = True; rm.font.name = "Calibri"
            rt = p.add_run(); rt.text = bt; rt.font.size = Pt(11)
            rt.font.color.rgb = accent_color; rt.font.bold = True; rt.font.name = "Calibri"
        else:
            rm = p.add_run(); rm.text = "\u2022  "; rm.font.size = Pt(11)
            rm.font.color.rgb = MED_GRAY; rm.font.name = "Calibri"
            rt = p.add_run(); rt.text = bt; rt.font.size = Pt(11)
            rt.font.color.rgb = DARK_TEXT; rt.font.bold = False; rt.font.name = "Calibri"


# ═══ SLIDE 7: Answer Key ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "ANSWER KEY", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.2), Inches(1.3), Inches(2.9), Inches(0.04), AWS_ORANGE)

answers = [
    ("A", "Standard + Glacier IR",    "CORRECT",   CORRECT_GRN,
     "Glacier Instant Retrieval saves 68% vs Standard-IA for once-per-quarter access."),
    ("B", "Standard + Standard-IA",   "INCORRECT", INCORRECT_R,
     "Works but not most cost-effective — Standard-IA costs 68% more than Glacier IR at quarterly access."),
    ("C", "Glacier IR + Deep Archive", "INCORRECT", INCORRECT_R,
     "Glacier IR wrong for frequent access; Deep Archive 12-48hr retrieval too slow for quarterly."),
    ("D", "Standard-IA + Glacier IR", "INCORRECT", INCORRECT_R,
     "Standard-IA wrong for frequent access — $0.01/GB retrieval fee adds up quickly."),
]

for i, (letter, name, verdict, color, reason) in enumerate(answers):
    y = Inches(1.8) + i * Inches(1.15)
    row_bg = RGBColor(0x2C, 0x3A, 0x4A) if i % 2 == 0 else RGBColor(0x34, 0x44, 0x55)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), row_bg)
    add_shape(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5), color)
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 letter, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(3.2), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)
    add_shape(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45), color)
    add_text_box(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45),
                 verdict, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.65),
                 reason, font_size=11, color=LIGHT_GRAY)

add_rect(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45), AWS_ORANGE)
add_text_box(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45),
             "Correct Answer:  A (S3 Standard + S3 Glacier Instant Retrieval)",
             font_size=15, color=AWS_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_ORANGE)

prs.save(OUT)
print(f"Saved: {OUT}")
