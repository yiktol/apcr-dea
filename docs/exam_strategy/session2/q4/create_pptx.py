#!/usr/bin/env python3
"""Generate PowerPoint for Q4 exam strategy."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "q4_exam_strategy.pptx")

AWS_DARK    = RGBColor(0x23, 0x2F, 0x3E)
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CORRECT_GRN = RGBColor(0x2E, 0x7D, 0x32)
INCORRECT_R = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY    = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT    = RGBColor(0x21, 0x21, 0x21)
SUBTLE_BG   = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    if border_color: s.line.color.rgb = border_color; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def add_rect(slide, left, top, width, height, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill_color; s.line.fill.background()
    return s

def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(font_size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font_name; p.alignment = alignment
    if line_spacing: p.line_spacing = Pt(line_spacing)
    return txBox


# ═══ SLIDE 1: Title ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), AWS_ORANGE)
add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
             "Database Migration Strategy",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
             "On-Premises to RDS  |  Schema Conversion  |  AWS SCT  |  AWS DMS",
             font_size=20, color=AWS_ORANGE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), AWS_ORANGE)
add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "Exam Strategy  —  Question 4  —  Most operationally efficient approach",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.43), SLIDE_W, Inches(0.07), AWS_ORANGE)


# ═══ SLIDE 2: Verbatim Question ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

add_shape(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75), AWS_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.65),
             "QUESTION", font_size=24, color=WHITE, bold=True)

scenario = (
    "A data engineer is planning to migrate their on-premises database to Amazon RDS. "
    "As part of the migration, they need to perform a schema conversion from the "
    "on-premises database to the Amazon RDS database. The data engineer also wants the "
    "solution to scan their application source code for embedded SQL statements and "
    "convert them as part of the project."
)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1),
             scenario, font_size=13, color=DARK_TEXT, line_spacing=20)

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.4),
             "What is the recommended approach to perform the schema conversion for the migration in the most operationally efficient way?",
             font_size=14, color=AWS_DARK, bold=True)

options_raw = [
    ("A.", "Use the AWS Database Migration Service (AWS DMS) to convert the database schema and migrate the data."),
    ("B.", "Manually convert the database schema using SQL scripts and then use AWS DMS to migrate the data."),
    ("C.", "Use the AWS Schema Conversion Tool (AWS SCT) to convert the database schema and then use AWS DMS to migrate the data."),
    ("D.", "Use the AWS Schema Conversion Tool (AWS SCT) to convert the database schema, and then use the AWS Transfer family to migrate the data the data."),
]

card_y = Inches(3.2)
card_h = Inches(0.72)
card_gap = Inches(0.08)

for i, (letter, text) in enumerate(options_raw):
    y = card_y + i * (card_h + card_gap)
    add_shape(slide, Inches(0.8), y, Inches(11.7), card_h, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45),
                 letter.replace(".", ""), font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(10.5), card_h - Inches(0.15),
                 text, font_size=12, color=DARK_TEXT)

# Timer GIF
timer_gif = os.path.join(os.path.dirname(BASE), "countdown_2min.gif")
if os.path.exists(timer_gif):
    slide.shapes.add_picture(timer_gif, Inches(10.5), Inches(6.5), Inches(2.5), Inches(0.85))


# ═══ SLIDES 3-6: Option details ═══
options_detail = [
    {
        "file": "option_a_dms_only.png",
        "letter": "A",
        "title": "AWS DMS Only",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("DMS migrates data between source and target databases (full load + CDC)", False),
            ("DMS does NOT perform schema conversion — it only moves data", True),
            ("Cannot scan application source code for embedded SQL statements", True),
            ("A separate tool (AWS SCT) is needed for schema and SQL conversion", False),
        ],
    },
    {
        "file": "option_b_manual_sql.png",
        "letter": "B",
        "title": "Manual SQL Scripts + DMS",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Manually writing DDL conversion scripts is possible but error-prone and slow", False),
            ("Not operationally efficient — AWS SCT automates this entire process", True),
            ("Cannot automatically scan application code for embedded SQL", True),
            ("DMS for data migration is correct, but the schema step adds unnecessary manual effort", False),
        ],
    },
    {
        "file": "option_c_sct_dms.png",
        "letter": "C",
        "title": "AWS SCT + AWS DMS",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("SCT auto-converts database schema: DDL, stored procedures, functions, views", False),
            ("SCT scans application source code for embedded SQL and converts it automatically", True),
            ("SCT generates an assessment report flagging items needing manual attention", False),
            ("DMS handles data migration with full load + change data capture (CDC)", False),
            ("Purpose-built combination: SCT for schema, DMS for data — most operationally efficient", True),
        ],
    },
    {
        "file": "option_d_sct_transfer.png",
        "letter": "D",
        "title": "AWS SCT + Transfer Family",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("SCT for schema conversion is correct — handles DDL and SQL conversion", False),
            ("AWS Transfer Family is for SFTP/FTPS/FTP file transfers to S3 — not database migration", True),
            ("Transfer Family cannot migrate database data to RDS or perform CDC replication", True),
            ("DMS is the correct service for migrating database data to Amazon RDS", False),
        ],
    },
]

for idx, opt in enumerate(options_detail):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)
    accent_color = CORRECT_GRN if opt["is_correct"] else INCORRECT_R

    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent_color)
    add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

    add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.75), AWS_DARK)
    add_shape(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55), accent_color)
    add_text_box(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55),
                 opt["letter"], font_size=24, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), Inches(0.38), Inches(8), Inches(0.6),
                 opt["title"], font_size=24, color=WHITE, bold=True)
    add_shape(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5), accent_color)
    add_text_box(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5),
                 opt["verdict"], font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Option text card
    option_text = options_raw[idx][1]
    oc_top = Inches(1.16)
    add_shape(slide, Inches(0.53), oc_top, Inches(12.27), Inches(0.65),
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_shape(slide, Inches(0.73), oc_top + Inches(0.1), Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(0.73), oc_top + Inches(0.1), Inches(0.45), Inches(0.45),
                 opt["letter"], font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.43), oc_top + Inches(0.1), Inches(11.17), Inches(0.30),
                 option_text, font_size=12, color=DARK_TEXT)

    # Diagram
    img_path = os.path.join(BASE, opt["file"])
    img = Image.open(img_path)
    aspect = img.size[0] / img.size[1]
    da_l, da_t, da_w, da_h = Inches(0.5), Inches(1.95), Inches(12.3), Inches(3.25)
    if da_w / da_h > aspect:
        fh = da_h; fw = int(da_h * aspect)
    else:
        fw = da_w; fh = int(da_w / aspect)
    slide.shapes.add_picture(img_path, da_l + (da_w - fw)//2, da_t + (da_h - fh)//2, fw, fh)

    # Explanation panel
    pt = Inches(5.36); ph = Inches(1.94)
    add_shape(slide, Inches(0.5), pt, Inches(12.3), ph, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))
    add_rect(slide, Inches(0.5), pt + Inches(0.19), Inches(0.06), Inches(1.65), accent_color)

    vl = "Correct" if opt["is_correct"] else "Incorrect"
    add_text_box(slide, Inches(0.85), pt + Inches(0.15), Inches(1.5), Inches(0.35),
                 vl, font_size=16, color=accent_color, bold=True)
    add_text_box(slide, Inches(2.2), pt + Inches(0.15), Inches(0.3), Inches(0.35),
                 "|", font_size=16, color=RGBColor(0xE0, 0xE0, 0xE0), bold=True)
    add_text_box(slide, Inches(2.45), pt + Inches(0.15), Inches(1.5), Inches(0.35),
                 "Explanation", font_size=14, color=AWS_DARK, bold=True)
    add_rect(slide, Inches(0.85), pt + Inches(0.55), Inches(11.7), Inches(0.015),
             RGBColor(0xE0, 0xE0, 0xE0))

    txBox = slide.shapes.add_textbox(Inches(0.85), pt + Inches(0.64), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame; tf.word_wrap = True
    for bi, (bt, is_key) in enumerate(opt["bullets"]):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.line_spacing = Pt(18); p.font.name = "Calibri"
        if is_key:
            rm = p.add_run(); rm.text = "\u25B6  "; rm.font.size = Pt(9)
            rm.font.color.rgb = accent_color; rm.font.bold = True; rm.font.name = "Calibri"
            rt = p.add_run(); rt.text = bt; rt.font.size = Pt(11)
            rt.font.color.rgb = accent_color; rt.font.bold = True; rt.font.name = "Calibri"
        else:
            rm = p.add_run(); rm.text = "\u2022  "; rm.font.size = Pt(11)
            rm.font.color.rgb = MED_GRAY; rm.font.name = "Calibri"
            rt = p.add_run(); rt.text = bt; rt.font.size = Pt(11)
            rt.font.color.rgb = DARK_TEXT; rt.font.bold = False; rt.font.name = "Calibri"


# ═══ SLIDE 7: Answer Key ═══
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "ANSWER KEY", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.2), Inches(1.3), Inches(2.9), Inches(0.04), AWS_ORANGE)

answers = [
    ("A", "DMS Only",              "INCORRECT", INCORRECT_R,
     "DMS migrates data only — does not convert schema or scan application SQL."),
    ("B", "Manual SQL + DMS",      "INCORRECT", INCORRECT_R,
     "Manual scripts are error-prone and slow — SCT automates the entire conversion."),
    ("C", "SCT + DMS",             "CORRECT",   CORRECT_GRN,
     "SCT converts schema + scans app SQL; DMS migrates data — purpose-built combination."),
    ("D", "SCT + Transfer Family", "INCORRECT", INCORRECT_R,
     "Transfer Family is for FTP file transfers to S3 — not for database data migration."),
]

for i, (letter, name, verdict, color, reason) in enumerate(answers):
    y = Inches(1.8) + i * Inches(1.15)
    row_bg = RGBColor(0x2C, 0x3A, 0x4A) if i % 2 == 0 else RGBColor(0x34, 0x44, 0x55)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), row_bg)
    add_shape(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5), color)
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 letter, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(3.2), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)
    add_shape(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45), color)
    add_text_box(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45),
                 verdict, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.65),
                 reason, font_size=11, color=LIGHT_GRAY)

add_rect(slide, Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.45), AWS_ORANGE)
add_text_box(slide, Inches(3.5), Inches(6.9), Inches(6.3), Inches(0.45),
             "Correct Answer:  C (AWS SCT + AWS DMS)",
             font_size=15, color=AWS_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_ORANGE)

prs.save(OUT)
print(f"Saved: {OUT}")
