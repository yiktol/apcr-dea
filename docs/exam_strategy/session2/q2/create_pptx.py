#!/usr/bin/env python3
"""
Generate PowerPoint for Q2 exam strategy.
- Slide 1: Title
- Slide 2: Verbatim question & options (no correct-answer highlighting)
- Slides 3-7: One per option (full-width diagram + horizontal bullet explanation)
- Slide 8: Answer key summary
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os

BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "q2_exam_strategy.pptx")

# ── Colors ──
AWS_DARK    = RGBColor(0x23, 0x2F, 0x3E)
AWS_ORANGE  = RGBColor(0xFF, 0x99, 0x00)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CORRECT_GRN = RGBColor(0x2E, 0x7D, 0x32)
INCORRECT_R = RGBColor(0xC6, 0x28, 0x28)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY    = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT    = RGBColor(0x21, 0x21, 0x21)
SUBTLE_BG   = RGBColor(0xFA, 0xFA, 0xFA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox


# ══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.07), AWS_ORANGE)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1),
             "Redshift Cost Optimization Strategy",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.7),
             "UNLOAD to S3  |  Parquet Format  |  Redshift Spectrum  |  Query Data Lake",
             font_size=20, color=AWS_ORANGE, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5), Inches(4.2), Inches(3.3), Inches(0.04), AWS_ORANGE)

add_text_box(slide, Inches(2), Inches(4.8), Inches(9.3), Inches(0.5),
             "Exam Strategy  —  Question 2  —  Choose TWO correct answers",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.43), SLIDE_W, Inches(0.07), AWS_ORANGE)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: Verbatim Question & Options
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)
add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

# Header
add_shape(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.75), AWS_DARK)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.65),
             "QUESTION", font_size=24, color=WHITE, bold=True)

# Scenario — verbatim
scenario = (
    "A data engineer is tasked with optimizing the cost of a Redshift cluster by "
    "relocating data that is older than 6 months into Amazon S3. The goal is to store "
    "this data in a way that minimizes expenses while still allowing it to be easily "
    "queried and combined with data in the cluster for reporting and analytics purposes."
)
add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.1),
             scenario, font_size=13, color=DARK_TEXT, line_spacing=20)

add_text_box(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.4),
             "Which combination of steps will provide an appropriate solution. (Choose two)",
             font_size=14, color=AWS_DARK, bold=True)

# Options — verbatim, uniform styling
options_raw = [
    ("A.", "Use the UNLOAD command to move data into S3 as JSON files"),
    ("B.", "Use Redshift Federated Query to access the data in S3."),
    ("C.", "Use the COPY command to move data into S3 as Parquet files."),
    ("D.", "Use Redshift Spectrum to access the data in S3"),
    ("E.", "Use the UNLOAD command to move data into S3 as Parquet files."),
]

card_y = Inches(3.2)
card_h = Inches(0.65)
card_gap = Inches(0.08)

for i, (letter, text) in enumerate(options_raw):
    y = card_y + i * (card_h + card_gap)

    add_shape(slide, Inches(0.8), y, Inches(11.7), card_h, SUBTLE_BG,
              border_color=RGBColor(0xE0, 0xE0, 0xE0))

    add_shape(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(1.0), y + Inches(0.1), Inches(0.45), Inches(0.45),
                 letter.replace(".", ""), font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.7), y + Inches(0.1), Inches(10.5), card_h - Inches(0.15),
                 text, font_size=12, color=DARK_TEXT)

# Timer GIF — bottom right of question slide
timer_gif = os.path.join(os.path.dirname(BASE), "countdown_2min.gif")
if os.path.exists(timer_gif):
    slide.shapes.add_picture(timer_gif, Inches(10.5), Inches(6.5), Inches(2.5), Inches(0.85))


# ═══ SLIDES 3-7: Option details ═══
options_detail = [
    {
        "file": "option_a_unload_json.png",
        "letter": "A",
        "title": "UNLOAD to S3 as JSON",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("UNLOAD is the correct command to export data from Redshift to S3", False),
            ("JSON is a row-based format — not optimized for analytics or columnar queries", True),
            ("Poor compression compared to Parquet — results in higher S3 storage costs", True),
            ("Full file scans required for queries — no column pruning or predicate pushdown", False),
            ("Parquet is the right format choice for cost-effective analytics on S3", False),
        ],
    },
    {
        "file": "option_b_federated_query.png",
        "letter": "B",
        "title": "Redshift Federated Query",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("Federated Query queries relational databases (RDS/Aurora PostgreSQL and MySQL)", False),
            ("Does NOT support querying data in Amazon S3 — wrong service for this use case", True),
            ("Redshift Spectrum is the correct feature for querying S3 data via SQL", True),
            ("Federated Query is designed for live joins with operational databases, not data lakes", False),
        ],
    },
    {
        "file": "option_c_copy_to_s3.png",
        "letter": "C",
        "title": "COPY Command to S3",
        "verdict": "INCORRECT",
        "is_correct": False,
        "bullets": [
            ("COPY loads data INTO Redshift from S3 — the direction is S3 to Redshift", True),
            ("Cannot be used to export data from Redshift to S3", True),
            ("UNLOAD is the correct command for exporting data OUT of Redshift to S3", False),
            ("Common exam trap: confusing COPY (inbound) with UNLOAD (outbound)", False),
        ],
    },
    {
        "file": "option_d_unload_parquet.png",
        "letter": "D",
        "title": "UNLOAD to S3 as Parquet",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("UNLOAD exports query results from Redshift to S3 — correct command direction", False),
            ("Parquet is columnar — optimized for analytics with column pruning and predicate pushdown", True),
            ("Built-in compression reduces S3 storage costs — key for the cost optimization goal", True),
            ("Supports PARTITION BY for organized, efficient data layout in S3", False),
            ("Self-describing schema embedded in files — no external schema management needed", False),
        ],
    },
    {
        "file": "option_e_spectrum.png",
        "letter": "E",
        "title": "Redshift Spectrum",
        "verdict": "CORRECT",
        "is_correct": True,
        "bullets": [
            ("Queries data directly in S3 using SQL — no need to load it back into Redshift", True),
            ("Allows JOINs between S3 external tables and local Redshift tables for reporting", True),
            ("Massively parallel Spectrum workers scan S3 data independently of cluster compute", False),
            ("Keeps historical data in cost-effective S3 while maintaining full query capability", False),
            ("Uses AWS Glue Data Catalog or Hive metastore for external table definitions", False),
        ],
    },
]

for idx, opt in enumerate(options_detail):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(slide, WHITE)

    accent_color = CORRECT_GRN if opt["is_correct"] else INCORRECT_R

    # Top accent
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent_color)
    # Bottom accent
    add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_DARK)

    # Header bar
    add_shape(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.75), AWS_DARK)

    # Letter badge in header
    add_shape(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55), accent_color)
    add_text_box(slide, Inches(0.7), Inches(0.38), Inches(0.55), Inches(0.55),
                 opt["letter"], font_size=24, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.5), Inches(0.38), Inches(8), Inches(0.6),
                 opt["title"], font_size=24, color=WHITE, bold=True)

    # Verdict badge
    add_shape(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5), accent_color)
    add_text_box(slide, Inches(10.5), Inches(0.42), Inches(2.1), Inches(0.5),
                 opt["verdict"], font_size=16, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # ── Option text card (below header, above diagram) ──
    option_text = options_raw[idx][1]
    opt_card_top = Inches(1.16)
    opt_card_h = Inches(0.65)

    add_shape(slide, Inches(0.53), opt_card_top, Inches(12.27), opt_card_h,
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))

    add_shape(slide, Inches(0.73), opt_card_top + Inches(0.1),
              Inches(0.45), Inches(0.45), MED_GRAY)
    add_text_box(slide, Inches(0.73), opt_card_top + Inches(0.1),
                 Inches(0.45), Inches(0.45),
                 opt["letter"], font_size=16, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.43), opt_card_top + Inches(0.1),
                 Inches(11.17), Inches(0.30),
                 option_text, font_size=12, color=DARK_TEXT)

    # ── Diagram: full-width, shifted down for option card ──
    img_path = os.path.join(BASE, opt["file"])
    img = Image.open(img_path)
    img_w, img_h = img.size
    aspect = img_w / img_h

    diagram_area_left = Inches(0.5)
    diagram_area_top = Inches(1.95)
    diagram_area_w = Inches(12.3)
    diagram_area_h = Inches(3.25)

    if diagram_area_w / diagram_area_h > aspect:
        final_h = diagram_area_h
        final_w = int(diagram_area_h * aspect)
    else:
        final_w = diagram_area_w
        final_h = int(diagram_area_w / aspect)

    img_left = diagram_area_left + (diagram_area_w - final_w) // 2
    img_top = diagram_area_top + (diagram_area_h - final_h) // 2

    slide.shapes.add_picture(img_path, img_left, img_top, final_w, final_h)

    # ── Explanation panel: horizontal bar at bottom (shifted up, taller) ──
    panel_top = Inches(5.36)
    panel_h = Inches(1.94)

    add_shape(slide, Inches(0.5), panel_top, Inches(12.3), panel_h,
              SUBTLE_BG, border_color=RGBColor(0xE0, 0xE0, 0xE0))

    # Accent stripe
    add_rect(slide, Inches(0.5), panel_top + Inches(0.19), Inches(0.06), Inches(1.65), accent_color)

    # Verdict label
    verdict_label = "Correct" if opt["is_correct"] else "Incorrect"
    add_text_box(slide, Inches(0.85), panel_top + Inches(0.15),
                 Inches(1.5), Inches(0.35),
                 verdict_label, font_size=16, color=accent_color, bold=True)

    # Separator
    add_text_box(slide, Inches(2.2), panel_top + Inches(0.15),
                 Inches(0.3), Inches(0.35),
                 "|", font_size=16, color=RGBColor(0xE0, 0xE0, 0xE0), bold=True)

    add_text_box(slide, Inches(2.45), panel_top + Inches(0.15),
                 Inches(1.5), Inches(0.35),
                 "Explanation", font_size=14, color=AWS_DARK, bold=True)

    # Thin line
    add_rect(slide, Inches(0.85), panel_top + Inches(0.55),
             Inches(11.7), Inches(0.015), RGBColor(0xE0, 0xE0, 0xE0))

    # Bullet-point explanation
    bullets = opt["bullets"]
    txBox = slide.shapes.add_textbox(
        Inches(0.85), panel_top + Inches(0.64), Inches(11.7), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True

    for bi, (bullet_text, is_key) in enumerate(bullets):
        p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
        p.line_spacing = Pt(18)
        p.font.name = "Calibri"

        if is_key:
            run_m = p.add_run()
            run_m.text = "\u25B6  "
            run_m.font.size = Pt(9)
            run_m.font.color.rgb = accent_color
            run_m.font.bold = True
            run_m.font.name = "Calibri"

            run_t = p.add_run()
            run_t.text = bullet_text
            run_t.font.size = Pt(11)
            run_t.font.color.rgb = accent_color
            run_t.font.bold = True
            run_t.font.name = "Calibri"
        else:
            run_m = p.add_run()
            run_m.text = "\u2022  "
            run_m.font.size = Pt(11)
            run_m.font.color.rgb = MED_GRAY
            run_m.font.name = "Calibri"

            run_t = p.add_run()
            run_t.text = bullet_text
            run_t.font.size = Pt(11)
            run_t.font.color.rgb = DARK_TEXT
            run_t.font.bold = False
            run_t.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════
# SLIDE 8: Answer Key Summary
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, AWS_DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), AWS_ORANGE)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.8),
             "ANSWER KEY", font_size=36, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(5.2), Inches(1.3), Inches(2.9), Inches(0.04), AWS_ORANGE)

answers = [
    ("A", "UNLOAD as JSON",        "INCORRECT", INCORRECT_R,
     "UNLOAD is correct command, but JSON is row-based — not optimized for analytics."),
    ("B", "Federated Query",       "INCORRECT", INCORRECT_R,
     "Queries RDS/Aurora databases, not S3 — wrong service for data lake access."),
    ("C", "COPY to S3",            "INCORRECT", INCORRECT_R,
     "COPY loads data INTO Redshift from S3 — cannot export data out to S3."),
    ("D", "UNLOAD as Parquet",     "CORRECT",   CORRECT_GRN,
     "UNLOAD exports to S3; Parquet is columnar, compressed — ideal for cost + analytics."),
    ("E", "Redshift Spectrum",     "CORRECT",   CORRECT_GRN,
     "Queries S3 data via SQL, JOINs with local tables — no reload needed."),
]

for i, (letter, name, verdict, color, reason) in enumerate(answers):
    y = Inches(1.8) + i * Inches(1.0)

    row_bg = RGBColor(0x2C, 0x3A, 0x4A) if i % 2 == 0 else RGBColor(0x34, 0x44, 0x55)
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), row_bg)

    add_shape(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5), color)
    add_text_box(slide, Inches(1.0), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 letter, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1.8), y + Inches(0.1), Inches(3.2), Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True)

    add_shape(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45), color)
    add_text_box(slide, Inches(5.2), y + Inches(0.18), Inches(1.7), Inches(0.45),
                 verdict, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(7.2), y + Inches(0.1), Inches(5.1), Inches(0.65),
                 reason, font_size=11, color=LIGHT_GRAY)

# Bottom callout
add_rect(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45), AWS_ORANGE)
add_text_box(slide, Inches(2.5), Inches(6.9), Inches(8.3), Inches(0.45),
             "Correct Answers:  D (UNLOAD as Parquet)  +  E (Redshift Spectrum)",
             font_size=15, color=AWS_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.44), SLIDE_W, Inches(0.06), AWS_ORANGE)

# ── Save ──
prs.save(OUT)
print(f"Saved: {OUT}")
