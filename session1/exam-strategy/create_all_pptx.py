"""
Generate PowerPoint presentations for each exam strategy question.
Uses python-pptx with improved layout and design.
Each PPT has:
  - Slide 1: Question + options + visually appealing 2-min timer
  - Slides 2+: One per option with diagram (not stretched), bullet explanation, option text on top
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

os.chdir(os.path.dirname(os.path.abspath(__file__)))

TIMER_PATH = "countdown_2min.gif"

# Slide dimensions (16:9)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

# Color palette
NAVY = RGBColor(0x0d, 0x1b, 0x2a)
DARK_NAVY = RGBColor(0x09, 0x12, 0x1c)
ACCENT_BLUE = RGBColor(0x1e, 0x88, 0xe5)
LIGHT_BLUE = RGBColor(0xbb, 0xde, 0xfb)
CORRECT_GREEN = RGBColor(0x00, 0xc8, 0x53)
INCORRECT_RED = RGBColor(0xff, 0x17, 0x44)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfa)
MID_GRAY = RGBColor(0x90, 0xa4, 0xae)
DARK_TEXT = RGBColor(0x26, 0x32, 0x38)
TIMER_BG = RGBColor(0x1a, 0x23, 0x7e)
TIMER_RING = RGBColor(0x00, 0xe6, 0x76)


def get_image_dimensions(path, max_width, max_height):
    """Get proportional dimensions for an image without stretching."""
    with Image.open(path) as img:
        w, h = img.size
    ratio = min(max_width / w, max_height / h)
    return int(w * ratio), int(h * ratio)


def set_slide_bg_gradient(slide, color1, color2):
    """Set gradient background."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[1].color.rgb = color2


def set_slide_bg_solid(slide, color):
    """Set solid background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, opacity=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def create_question_slide(prs, title, question_text, options):
    """Create slide 1: question + options + timer."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_gradient(slide, NAVY, DARK_NAVY)

    # Left accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(6.5), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Segoe UI"

    # Question text
    q_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(7.2), Inches(1.5))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question_text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"
    p.space_after = Pt(4)

    # Options with letter badges
    opt_top = 2.4
    for i, opt in enumerate(options):
        letter = chr(65 + i)

        # Letter badge circle
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(opt_top - 0.02), Inches(0.32), Inches(0.32)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.margin_top = Pt(0)
        badge_tf.margin_bottom = Pt(0)
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = badge_tf.paragraphs[0]
        bp.text = letter
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER

        # Option text
        opt_box = slide.shapes.add_textbox(Inches(0.95), Inches(opt_top), Inches(7.0), Inches(0.4))
        tf = opt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = opt
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xec, 0xef, 0xf1)
        p.font.name = "Segoe UI"

        opt_top += 0.58

    # Timer area (right side) - visually appealing card
    timer_card = add_rounded_rect(
        slide, Inches(7.7), Inches(0.3), Inches(2.1), Inches(2.3), TIMER_BG
    )

    # Timer label
    timer_label = slide.shapes.add_textbox(Inches(7.8), Inches(0.4), Inches(1.9), Inches(0.35))
    tf = timer_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "⏱ TIME"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TIMER_RING
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    # Timer GIF
    if os.path.exists(TIMER_PATH):
        slide.shapes.add_picture(TIMER_PATH, Inches(8.0), Inches(0.7), Inches(1.7), Inches(1.7))

    return slide


def create_option_slide(prs, letter, option_text, diagram_path, explanation_bullets, is_correct):
    """Create a slide for one option with diagram and bullet explanation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, LIGHT_GRAY)

    # Status color
    status_color = CORRECT_GREEN if is_correct else INCORRECT_RED
    status_text = "✓ CORRECT" if is_correct else "✗ INCORRECT"

    # Top banner
    banner = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.9)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    # Status badge (left)
    badge = add_rounded_rect(slide, Inches(0.3), Inches(0.2), Inches(1.4), Inches(0.45), status_color)
    badge_tf = badge.text_frame
    badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_tf.margin_left = Pt(4)
    badge_tf.margin_right = Pt(4)
    bp = badge_tf.paragraphs[0]
    bp.text = status_text
    bp.font.size = Pt(10)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.font.name = "Segoe UI"
    bp.alignment = PP_ALIGN.CENTER

    # Option text on banner
    opt_box = slide.shapes.add_textbox(Inches(1.9), Inches(0.15), Inches(7.8), Inches(0.65))
    tf = opt_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"Option {letter}: {option_text}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    # Diagram (left side, proportional - not stretched)
    if os.path.exists(diagram_path):
        max_w_px = 580  # max width in pixels for display area
        max_h_px = 340  # max height in pixels for display area
        max_w_in = 5.8
        max_h_in = 3.4

        with Image.open(diagram_path) as img:
            img_w, img_h = img.size

        # Calculate proportional size
        ratio = min(max_w_in / (img_w / 96.0), max_h_in / (img_h / 96.0))
        display_w = (img_w / 96.0) * ratio
        display_h = (img_h / 96.0) * ratio

        # Cap to max
        if display_w > max_w_in:
            display_w = max_w_in
            display_h = (img_h / img_w) * display_w
        if display_h > max_h_in:
            display_h = max_h_in
            display_w = (img_w / img_h) * display_h

        # Center vertically in available space
        img_top = 1.1 + (3.6 - display_h) / 2
        slide.shapes.add_picture(diagram_path, Inches(0.2), Inches(img_top),
                                 Inches(display_w), Inches(display_h))

    # Explanation panel (right side) - card style
    panel_left = Inches(6.2)
    panel_top = Inches(1.1)
    panel_w = Inches(3.6)
    panel_h = Inches(4.2)

    panel = add_rounded_rect(slide, panel_left, panel_top, panel_w, panel_h, WHITE)

    # Accent line on panel
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(1.1), Inches(0.06), panel_h
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = status_color
    accent_line.line.fill.background()

    # "Why?" header
    why_box = slide.shapes.add_textbox(Inches(6.45), Inches(1.25), Inches(3.2), Inches(0.35))
    tf = why_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Why this is correct:" if is_correct else "Why this is incorrect:"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = status_color
    p.font.name = "Segoe UI"

    # Bullet points
    bullets_box = slide.shapes.add_textbox(Inches(6.45), Inches(1.65), Inches(3.2), Inches(3.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(explanation_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)

    return slide


# ============================================================
# QUESTION DATA
# ============================================================

questions = [
    {
        "folder": "q1",
        "title": "Kinesis Throughput Exception",
        "question": "An Amazon Kinesis application is trying to read data from a Kinesis data stream. However, the read data call is rejected. The following error message is displayed: ProvisionedThroughputExceededException.\n\nWhich combination of steps will resolve the error? (Select TWO)",
        "options": [
            "Configure enhanced fan-out on the stream",
            "Enable enhanced monitoring on the stream",
            "Increase the size of the GetRecords requests",
            "Increase the number of shards within the stream to provide enough capacity for the read data calls",
            "Make the application retry to read data from the stream",
        ],
        "diagrams": [
            "option_a_enhanced_fanout.png",
            "option_b_enhanced_monitoring.png",
            "option_c_larger_getrecords.png",
            "option_d_more_shards.png",
            "option_e_retry.png",
        ],
        "correct": [3, 4],
        "explanations": [
            ["Enhanced fan-out gives dedicated 2MB/s per consumer via HTTP/2",
             "Doesn't fix the root cause: shard-level read limit exceeded",
             "Only helps when multiple consumers compete for the same shard"],
            ["Provides shard-level CloudWatch metrics",
             "Observability tool only — doesn't increase throughput",
             "Useful for diagnosis, not resolution"],
            ["Larger batch size retrieves more records per call",
             "Does NOT increase the 2 MB/s per-shard read limit",
             "May actually worsen throttling by consuming more bandwidth per request"],
            ["Each shard provides 2 MB/s read throughput",
             "More shards = higher aggregate read capacity",
             "Directly addresses the throughput bottleneck"],
            ["Exponential backoff handles transient throttling",
             "Exception may be temporary due to burst traffic",
             "AWS SDK best practice for handling throughput errors"],
        ],
    },
    {
        "folder": "q2",
        "title": "Data Pipeline Architecture",
        "question": "A data engineer at a multinational corporation is designing a new data pipeline architecture. The company needs to consolidate data from multiple transactional databases, enable data analysis for business intelligence teams, and ensure the production databases' performance isn't impacted.\n\nWhat is the recommended approach?",
        "options": [
            "Direct all analytics queries to the transactional databases to ensure real-time data access",
            "Land data in a data lake, prepare it, then move selected data to a data warehouse for reporting",
            "Store all data directly in the data warehouse and perform transformations during query runtime",
            "Create separate data marts for each business unit without a centralized data warehouse",
        ],
        "diagrams": [
            "option_a_direct_queries.png",
            "option_b_lake_warehouse.png",
            "option_c_direct_warehouse.png",
            "option_d_data_marts.png",
        ],
        "correct": [1],
        "explanations": [
            ["Complex analytics queries degrade production DB performance",
             "Violates the requirement to protect production workloads",
             "No separation between OLTP and OLAP concerns"],
            ["Separates analytics from production databases",
             "S3 data lake provides flexible, scalable raw storage",
             "Glue transforms data; Redshift serves fast BI queries",
             "AWS recommended lake-to-warehouse pattern"],
            ["No data lake stage limits flexibility",
             "Query-time transforms are expensive and slow",
             "Raw data in warehouse increases storage costs"],
            ["Creates data silos across business units",
             "Duplicate pipelines increase operational overhead",
             "Cross-unit analysis becomes impossible"],
        ],
    },
    {
        "folder": "q3",
        "title": "PII Detection & Masking",
        "question": "A company is collecting data in an Amazon S3 data lake that includes personally identifiable information (PII). They need an automated solution to identify PII in new and existing data, provide an overview of findings, and invoke a masking application in real time when PII is found.\n\nWhich solution meets these requirements with the LEAST operational overhead?",
        "options": [
            "Create an AWS Lambda function to analyze data for PII. Configure S3 bucket notifications to invoke the Lambda when a new object is uploaded.",
            "Configure S3 bucket notifications with an Amazon EventBridge rule for new object uploads. Set the masking application as the target.",
            "Enable Amazon Macie. Create an EventBridge rule for Macie findings. Set the masking application as the target.",
            "Enable Amazon Macie. Create a Lambda function on a schedule to poll Macie findings and invoke the masking application.",
        ],
        "diagrams": [
            "option_a_custom_lambda.png",
            "option_b_eventbridge_upload.png",
            "option_c_macie_eventbridge.png",
            "option_d_macie_polling.png",
        ],
        "correct": [2],
        "explanations": [
            ["Requires building and maintaining custom PII detection logic",
             "Only scans new uploads, not existing data",
             "No built-in reporting or overview of findings",
             "High operational overhead"],
            ["No PII identification step at all",
             "Triggers masking on every upload blindly",
             "Wasteful — processes files without PII",
             "Doesn't meet the 'identify PII' requirement"],
            ["Macie is purpose-built for PII discovery in S3",
             "Scans both new and existing data automatically",
             "EventBridge provides real-time event-driven response",
             "Fully managed with built-in reporting dashboard"],
            ["Polling adds latency — not real-time",
             "Scheduled Lambda adds operational overhead",
             "May miss findings between poll intervals",
             "Event-driven approach is simpler and faster"],
        ],
    },
    {
        "folder": "q4",
        "title": "NFS Migration for Lambda",
        "question": "A company has data in an on-premises NFS file share and plans to migrate to AWS. They use AWS Lambda functions to analyze the data and want to continue using NFS. The data must be shared across all concurrently running Lambda functions.\n\nWhich solution should the company use?",
        "options": [
            "Migrate the data into the local storage for each Lambda function. Use the local storage for data access.",
            "Migrate the data to Amazon Elastic Block Store (Amazon EBS) volumes. Access the EBS volumes from the Lambda functions.",
            "Migrate the data to Amazon DynamoDB. Ensure the Lambda functions have permissions to access the table.",
            "Migrate the data to Amazon Elastic File System (Amazon EFS). Configure the Lambda functions to mount the file system.",
        ],
        "diagrams": [
            "option_a_local_storage.png",
            "option_b_ebs.png",
            "option_c_dynamodb.png",
            "option_d_efs.png",
        ],
        "correct": [3],
        "explanations": [
            ["Lambda /tmp is ephemeral (max 10 GB)",
             "Isolated per invocation — no sharing possible",
             "Data lost when function execution ends"],
            ["EBS can only attach to EC2 instances",
             "Lambda does not support EBS volume mounting",
             "Architecturally incompatible"],
            ["DynamoDB is a key-value database, not a file system",
             "Not NFS compatible — violates the requirement",
             "Would require complete application rewrite"],
            ["EFS provides NFS v4 compatible shared storage",
             "Lambda can mount EFS via VPC configuration",
             "All concurrent invocations access the same files",
             "Meets both NFS and shared access requirements"],
        ],
    },
    {
        "folder": "q5",
        "title": "VPC Security for Data Pipeline",
        "question": "A data engineer is designing a data processing pipeline that ingests sensitive customer data, processes it using EC2 instances in a VPC, and stores results in Amazon S3. The EC2 instances cannot be directly accessed from the internet.\n\nWhich combination of VPC security features should be implemented?",
        "options": [
            "Use public subnets for EC2 instances and rely solely on security groups to control access",
            "Place EC2 instances in private subnets, use a NAT gateway for internet access, and implement security groups",
            "Use VPC endpoints for S3 access and place EC2 instances in public subnets with restrictive NACLs",
            "Implement VPC peering between multiple VPCs, each containing a different part of the pipeline",
        ],
        "diagrams": [
            "option_a_public_sg.png",
            "option_b_private_nat.png",
            "option_c_vpc_endpoints_public.png",
            "option_d_vpc_peering.png",
        ],
        "correct": [1],
        "explanations": [
            ["Public subnets assign public IPs to instances",
             "Instances are directly reachable from the internet",
             "Violates the core security requirement",
             "Security groups alone are not defense-in-depth"],
            ["Private subnets have no internet gateway route",
             "NAT gateway enables outbound-only connectivity",
             "Security groups control instance-level traffic",
             "Meets all stated security requirements"],
            ["VPC endpoints for S3 are good practice",
             "BUT public subnets still expose instances to internet",
             "NACLs are stateless and harder to manage",
             "Doesn't meet the 'no internet access' requirement"],
            ["VPC peering enables inter-VPC communication",
             "Doesn't address internet access prevention",
             "Adds unnecessary architectural complexity",
             "Over-engineered for a single-pipeline use case"],
        ],
    },
]


# ============================================================
# GENERATE PRESENTATIONS
# ============================================================

for q in questions:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: Question + Options + Timer
    create_question_slide(prs, q["title"], q["question"], q["options"])

    # Option slides
    for i, (opt, diag, expl) in enumerate(zip(q["options"], q["diagrams"], q["explanations"])):
        letter = chr(65 + i)
        diagram_path = os.path.join(q["folder"], diag)
        is_correct = i in q["correct"]
        create_option_slide(prs, letter, opt, diagram_path, expl, is_correct)

    # Save
    output_path = os.path.join(q["folder"], f"{q['folder']}_exam_strategy.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")

print("\nAll presentations generated successfully!")
