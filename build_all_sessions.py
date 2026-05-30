"""
Build exam-strategy materials for sessions 2-5:
  1. Create markdown files (q1.md, q2.md, etc.) per question
  2. Generate architecture diagrams using the diagrams library
  3. Generate PowerPoint presentations with timer, diagrams, and explanations
"""
import os
import sys
from pathlib import Path

# Session data definitions
SESSIONS = {}

# ============================================================
# SESSION 2 QUESTIONS
# ============================================================
SESSIONS[2] = {
    "base_dir": "session2/exam-strategy",
    "questions": [
        {
            "id": "q1",
            "title": "Data Transfer to S3",
            "question": "Your customer is an enterprise sized business in the telecommunications industry. They are looking to setup a hybrid environment and intend to utilize AWS Redshift for querying a dataset comprising 5 PB of historical cell phone usage data. They have an ongoing data accumulation rate of 3 TB per month. The requirement is for the new incoming monthly data to be accessible for querying with less than a two-day delay.\n\nWhich combination of steps will provide an appropriate solution? (Select TWO)",
            "options": [
                "Leverage AWS Direct Connect for both the initial bulk transfer of historical data and the ongoing monthly data transfers, ensuring consistent network performance.",
                "Implement AWS DataSync for continuous synchronization of the 3 TB monthly data to Amazon S3, maintaining the required freshness of less than two days.",
                "Configure Amazon S3 Transfer Acceleration for the ongoing 3 TB monthly data uploads to minimize transfer times and meet the delay requirements.",
                "Use AWS Snowball for the initial migration of the 5 PB of historical data to Amazon S3, ensuring a fast and secure data transfer method for large datasets.",
                "Directly upload the ongoing 3 TB of data each month via the Internet using the AWS CLI or SDKs, relying on your existing internet bandwidth.",
            ],
            "correct": [1, 3],
            "explanations": [
                ["Dedicated 1-400 Gbps private network connection", "Not optimized for PB-scale bulk migration", "Snowball and DataSync are purpose-built for these requirements", "Requires physical cross-connect setup adding weeks of lead time"],
                ["Transfers data from on-premises to S3 with built-in scheduling", "Purpose-built protocol delivers up to 10 Gbps per task", "Meets the < 2-day freshness requirement", "Includes in-transit and at-rest integrity validation"],
                ["Speeds uploads via CloudFront edge locations", "No built-in scheduling or incremental sync", "Must be managed manually", "DataSync is better for continuous synchronization"],
                ["Physical storage devices (~210 TB each) shipped to data center", "Optimized for large one-time migrations — 5 PB is perfect", "Physical transfer bypasses all network bottlenecks", "256-bit encryption with tamper-resistant enclosure"],
                ["Relies entirely on existing internet bandwidth", "No transfer optimization or automatic retry", "At 100 Mbps, 3 TB takes ~2.8 days — exceeds SLA", "DataSync provides scheduling and optimized throughput natively"],
            ],
        },
        {
            "id": "q2",
            "title": "Redshift Data Offloading",
            "question": "A data engineer is tasked with optimizing the cost of a Redshift cluster by relocating data that is older than 6 months into Amazon S3. The goal is to store this data in a way that minimizes expenses while still allowing it to be easily queried and combined with data in the cluster for reporting and analytics purposes.\n\nWhich combination of steps will provide an appropriate solution? (Choose two)",
            "options": [
                "Use the UNLOAD command to move data into S3 as JSON files",
                "Use Redshift Federated Query to access the data in S3",
                "Use the COPY command to move data into S3 as Parquet files",
                "Use Redshift Spectrum to access the data in S3",
                "Use the UNLOAD command to move data into S3 as Parquet files",
            ],
            "correct": [3, 4],
            "explanations": [
                ["UNLOAD is the correct command to export from Redshift", "JSON is row-based — not optimized for analytics", "Poor compression vs Parquet — higher storage costs", "No column pruning or predicate pushdown"],
                ["Federated Query queries RDS/Aurora databases", "Does NOT support querying data in S3", "Redshift Spectrum is the correct feature for S3", "Designed for live joins with operational DBs"],
                ["COPY loads data INTO Redshift from S3", "Cannot export data from Redshift to S3", "Direction is wrong: S3 → Redshift only", "Common exam trap: COPY (inbound) vs UNLOAD (outbound)"],
                ["Queries data directly in S3 using SQL", "Allows JOINs between S3 external tables and local tables", "Massively parallel Spectrum workers scan S3 independently", "Uses Glue Data Catalog for external table definitions"],
                ["UNLOAD exports query results from Redshift to S3", "Parquet is columnar — optimized for analytics", "Built-in compression reduces S3 storage costs", "Supports PARTITION BY for organized data layout"],
            ],
        },
        {
            "id": "q3",
            "title": "Data Lake Integration",
            "question": "A financial services company maintains a vast and secure data lake built on Amazon S3. They analyze this data using Apache Hive, Amazon Athena, Amazon Redshift, and Amazon QuickSight. The company needs to identify a solution to integrate the data lake and the tools used by the engineering teams.\n\nWhat solution accomplishes this with the LEAST amount of operational overhead?",
            "options": [
                "Use AWS Glue to catalog the data and make it accessible to all the analysis tools.",
                "Directly connect each analysis tool to Amazon S3, configuring each tool individually.",
                "Export data from Amazon S3 to an Amazon EC2 instance and use it as a central analysis server.",
                "Upload copies of the data to each analysis tool separately.",
            ],
            "correct": [0],
            "explanations": [
                ["Glue crawlers auto-discover schema from S3", "Single catalog shared by Athena, Spectrum, Hive, QuickSight", "Fully serverless — no infrastructure to manage", "Schema changes propagate automatically to all tools"],
                ["Each tool must define its own schema independently", "Significant operational overhead maintaining schemas", "No shared metadata — risk of schema drift", "Schema changes require updating all four tools manually"],
                ["Requires managing EC2: OS patching, scaling, backups", "Adds unnecessary infrastructure overhead", "Data must be copied from S3 — transfer cost and staleness", "Glue provides same benefit as fully managed service"],
                ["Creates 4 separate copies — massive storage cost", "Data duplication leads to inconsistencies", "Must maintain 4 separate sync pipelines", "Highest operational overhead of all options"],
            ],
        },
        {
            "id": "q4",
            "title": "Schema Conversion & Migration",
            "question": "A data engineer is planning to migrate their on-premises database to Amazon RDS. They need to perform a schema conversion and also want the solution to scan application source code for embedded SQL statements and convert them.\n\nWhat is the recommended approach to perform the schema conversion in the most operationally efficient way?",
            "options": [
                "Use AWS Database Migration Service (AWS DMS) to convert the database schema and migrate the data.",
                "Manually convert the database schema using SQL scripts and then use AWS DMS to migrate the data.",
                "Use the AWS Schema Conversion Tool (AWS SCT) to convert the database schema and then use AWS DMS to migrate the data.",
                "Use the AWS Schema Conversion Tool (AWS SCT) to convert the database schema, and then use the AWS Transfer Family to migrate the data.",
            ],
            "correct": [2],
            "explanations": [
                ["DMS migrates data between databases (full load + CDC)", "DMS does NOT perform schema conversion", "Cannot scan application source code for SQL", "A separate tool (SCT) is needed for schema conversion"],
                ["Manual DDL conversion is error-prone and slow", "Not operationally efficient — SCT automates this", "Cannot automatically scan app code for embedded SQL", "DMS for data is correct but schema step is manual"],
                ["SCT auto-converts schema: DDL, stored procedures, views", "SCT scans application source code for embedded SQL", "SCT generates assessment report for manual items", "DMS handles data migration with full load + CDC"],
                ["SCT for schema conversion is correct", "AWS Transfer Family is for SFTP/FTP file transfers to S3", "Transfer Family cannot migrate database data to RDS", "DMS is the correct service for database data migration"],
            ],
        },
        {
            "id": "q5",
            "title": "S3 Storage Classes",
            "question": "A video streaming service stores large amounts of video content on Amazon S3. They have frequently accessed content that needs to be readily available and older content accessed once per quarter. The data engineer is tasked with optimizing storage costs while maintaining required performance.\n\nWhich combination of Amazon S3 storage classes will meet the requirements in the most cost effective way?",
            "options": [
                "S3 Standard for frequently accessed content and S3 Glacier Instant Retrieval for infrequently accessed content.",
                "S3 Standard for frequently accessed content and S3 Infrequent-Access for infrequently accessed content.",
                "S3 Glacier Instant Retrieval for frequently accessed content and S3 Glacier Deep Archive for infrequently accessed archive content.",
                "S3 Standard-IA for frequently accessed content, S3 Glacier Instant Retrieval for infrequently accessed archive content.",
            ],
            "correct": [0],
            "explanations": [
                ["S3 Standard for hot content — no retrieval fee", "Glacier Instant Retrieval saves 68% vs Standard-IA at quarterly access", "Still provides millisecond access — no latency trade-off", "AWS recommends this for data accessed once per quarter"],
                ["S3 Standard for hot content is correct", "Standard-IA costs 68% more than Glacier IR at quarterly access", "Standard-IA better suited for monthly access patterns", "Functional but not the MOST cost-effective option"],
                ["Glacier IR NOT designed for frequently accessed content", "High retrieval costs on frequent access", "Deep Archive has 12-48 hour retrieval — too slow", "Frequently accessed content requires S3 Standard"],
                ["Standard-IA NOT suitable for frequent access", "$0.01/GB retrieval fee accumulates quickly", "Glacier IR for quarterly content is correct", "But the hot tier choice is wrong — needs S3 Standard"],
            ],
        },
    ],
}

# ============================================================
# SESSION 3 QUESTIONS
# ============================================================
SESSIONS[3] = {
    "base_dir": "session3/exam-strategy",
    "questions": [
        {
            "id": "q1",
            "title": "Athena Query Optimization",
            "question": "A sports organization stores extensive web server logs on Amazon S3 from historic sporting events and intends to analyze this data swiftly using Amazon Athena. The majority of their queries are operational, focusing on logs from a specific day or sporting match.\n\nHow should the organization structure their log data in the MOST cost effective manner?",
            "options": [
                "Keep logs in their original format and rely on Amazon Athena's automatic data partitioning.",
                "Partition logs by year, then by month, and use gzip compression for each file.",
                "Store all logs in a single, large file and use Athena's built-in query optimization features.",
                "Partition logs by day and convert them into Apache Parquet format.",
            ],
            "correct": [3],
            "explanations": [
                ["Athena does NOT automatically partition data", "Full table scans on raw format are expensive", "No compression or columnar optimization", "Highest cost per query of all options"],
                ["Year/month partitioning is too coarse for daily queries", "Gzip is splittable but row-based — no column pruning", "Queries for a specific day still scan entire month", "Parquet + day partition is more cost-effective"],
                ["Single large file forces full scan on every query", "No partition pruning possible", "Athena charges per TB scanned — maximum cost", "Worst possible approach for cost optimization"],
                ["Day-level partitioning matches query patterns exactly", "Parquet is columnar — enables column pruning", "Parquet has built-in compression (snappy/zstd)", "Minimizes data scanned = minimizes Athena cost"],
            ],
        },
        {
            "id": "q2",
            "title": "EMR Cost Optimization",
            "question": "A company recently migrated a batch job to an Amazon EMR cluster that uses On-Demand Instances. The job is critical to operations with an SLA of 3 hours and takes an average of 2 hours to complete. The company wants to reduce costs while minimizing impact on availability.\n\nWhich combination of steps should a data analytics specialist take? (Select TWO)",
            "options": [
                "Configure Amazon EMR to use instance fleets with a provisioning timeout for the core nodes.",
                "Configure Amazon EMR uniform instance groups.",
                "Change the instance purchasing option for core nodes from On-Demand Instances to Spot Instances.",
                "Use Spot Instances for the task nodes.",
                "Use Spot Instances for all node types in the cluster.",
            ],
            "correct": [0, 3],
            "explanations": [
                ["Instance fleets allow mixing instance types for availability", "Provisioning timeout ensures capacity is obtained within SLA", "Diversifies across instance types to reduce interruption risk", "Best practice for critical workloads needing cost savings"],
                ["Uniform instance groups use single instance type", "Less flexible — higher risk of capacity unavailability", "No provisioning timeout capability", "Instance fleets are the better choice for cost optimization"],
                ["Core nodes store HDFS data — interruption causes data loss", "Spot interruption on core nodes can fail the entire job", "Critical workloads should keep core nodes on On-Demand", "Risk to availability is too high for a 3-hour SLA job"],
                ["Task nodes only process data — no HDFS storage", "Spot interruption on task nodes doesn't lose data", "Significant cost savings with minimal availability risk", "Job continues on remaining nodes if Spot is reclaimed"],
                ["Master node on Spot = cluster failure if interrupted", "Core nodes on Spot = data loss risk", "Too risky for a critical job with 3-hour SLA", "Only task nodes are safe for Spot in critical workloads"],
            ],
        },
        {
            "id": "q3",
            "title": "Glue Workflow Orchestration",
            "question": "A transportation company uses AWS Glue ETL jobs to run transformations on data lake data. A data engineer must visualize dependencies between transformations and run multiple transformations concurrently to reduce processing time.\n\nHow can the data engineer meet these requirements with the LEAST operational overhead?",
            "options": [
                "Use a scheduled trigger to start an AWS Glue workflow that will launch the required AWS Glue jobs.",
                "Use Amazon EventBridge to schedule an event to initiate an AWS Lambda function that will launch the required AWS Glue jobs.",
                "Use Amazon EventBridge to schedule an event to initiate an AWS Step Functions workflow that will launch the required AWS Glue jobs.",
                "Launch an Amazon EC2 instance that runs Apache Airflow to start the required AWS Glue jobs.",
            ],
            "correct": [0],
            "explanations": [
                ["Glue workflows natively visualize job dependencies", "Built-in concurrent execution of independent jobs", "Scheduled triggers are fully managed — no extra services", "Least operational overhead — single service solution"],
                ["Lambda can start Glue jobs but no dependency visualization", "Must build custom orchestration logic in Lambda", "Additional service to manage (EventBridge + Lambda)", "More overhead than native Glue workflows"],
                ["Step Functions can orchestrate but adds complexity", "Two additional services (EventBridge + Step Functions)", "Glue workflows already provide this natively", "Over-engineered for Glue-to-Glue orchestration"],
                ["EC2 + Airflow requires infrastructure management", "OS patching, scaling, Airflow maintenance", "Highest operational overhead of all options", "Managed alternatives exist for this use case"],
            ],
        },
        {
            "id": "q4",
            "title": "EMR Idle Cluster Termination",
            "question": "A company uses Amazon EMR clusters for workload analysis. Data engineers often forget to terminate clusters after work is complete, causing unnecessary costs. The company needs an automated solution to terminate idle clusters.\n\nWhich solution meets these requirements?",
            "options": [
                "Tag EMR clusters. Create a CloudWatch alarm on estimated monthly bill. When billing exceeds threshold, invoke a Lambda function to terminate clusters.",
                "Configure a CloudWatch metrics alarm on the IsIdle metric from EMR clusters to publish to an SNS topic. Subscribe a Lambda function to terminate the clusters.",
                "Implement a bash script on the primary node scheduled every 5 minutes to terminate clusters after 8 hours each day.",
                "Use AWS Trusted Advisor cost optimization to check for idle EMR clusters. Terminate using the AWS CLI.",
            ],
            "correct": [1],
            "explanations": [
                ["Billing alarms react to cost, not idle state", "Cluster may be idle long before billing threshold", "Not targeted — terminates based on cost not usage", "Delayed response — damage already done to budget"],
                ["IsIdle is a native EMR CloudWatch metric", "Directly detects when cluster has no running jobs", "SNS + Lambda provides automated real-time response", "Purpose-built metric for exactly this use case"],
                ["Hardcoded 8-hour limit is arbitrary and inflexible", "Bash script on primary node is fragile", "Doesn't detect actual idle state", "Not a scalable or maintainable solution"],
                ["Trusted Advisor provides recommendations, not automation", "Requires manual intervention to act on findings", "Not real-time — periodic checks only", "Doesn't meet the 'automated solution' requirement"],
            ],
        },
        {
            "id": "q5",
            "title": "Redshift ML",
            "question": "A retail company uses Amazon Redshift for transactional data. They want to leverage machine learning models for insights. Data preparation and feature engineering are time-consuming. They need an efficient solution leveraging existing Redshift data.\n\nWhich approach should you take?",
            "options": [
                "Use Amazon Redshift ML to create and train machine learning models using SQL directly within Amazon Redshift.",
                "Extract data from Redshift into S3, then use Amazon SageMaker to build, train, and deploy models.",
                "Provision EC2 instances with Apache Spark, extract data from Redshift, and perform data preparation before training models.",
                "Use AWS Glue to extract data from Redshift, perform transformation with Spark, then load into SageMaker for training.",
            ],
            "correct": [0],
            "explanations": [
                ["Redshift ML uses SQL — no data movement needed", "Automatically handles feature engineering", "Integrates with SageMaker Autopilot behind the scenes", "Most efficient — leverages existing data in place"],
                ["Requires data extraction step to S3", "Additional data movement and storage costs", "More complex pipeline to manage", "Works but not the most efficient approach"],
                ["Requires provisioning and managing EC2 infrastructure", "Manual data extraction and preparation", "Highest operational overhead", "Self-managed Spark adds significant complexity"],
                ["Multiple services to orchestrate (Glue + SageMaker)", "Data must be extracted and transformed externally", "More operational overhead than Redshift ML", "Over-engineered when Redshift ML exists"],
            ],
        },
    ],
}

# ============================================================
# SESSION 4 QUESTIONS
# ============================================================
SESSIONS[4] = {
    "base_dir": "session4/exam-strategy",
    "questions": [
        {
            "id": "q1",
            "title": "S3 Upload Integrity",
            "question": "An international investment banking firm frequently uploads 100 MB files to Amazon S3 for analysis with Amazon Athena. Recently, some uploads have become corrupted. They need to ensure successful upload and verify file contents remain unchanged from on-premises to S3.\n\nWhich recommendation is the MOST cost effective?",
            "options": [
                "Use the AWS CLI to upload files with the --sse parameter for server-side encryption, ensuring data integrity during transit.",
                "Utilize the S3 Multipart upload feature with MD5 checksum validation for each part, providing a strong guarantee that files are uploaded successfully and remain unaltered.",
                "Enable versioning on the S3 bucket to keep multiple versions of an object, allowing the company to revert to uncorrupted versions if necessary.",
                "Implement S3 Transfer Acceleration to increase the speed of file transfers and reduce the likelihood of corruption due to connection issues.",
            ],
            "correct": [1],
            "explanations": [
                ["SSE encrypts data at rest, not integrity validation", "Encryption doesn't detect or prevent upload corruption", "Doesn't verify content matches source", "Solves a different problem (confidentiality, not integrity)"],
                ["Multipart upload splits file into verifiable parts", "MD5 checksum validates each part independently", "S3 rejects parts that fail checksum — guarantees integrity", "100 MB files benefit from parallel part uploads"],
                ["Versioning keeps old copies but doesn't prevent corruption", "Doesn't detect corruption at upload time", "Reactive approach — corruption already happened", "Adds storage cost without solving root cause"],
                ["Transfer Acceleration speeds up uploads via edge locations", "Speed doesn't guarantee data integrity", "Doesn't validate content matches source", "Addresses latency, not corruption detection"],
            ],
        },
        {
            "id": "q2",
            "title": "S3 Object Lambda",
            "question": "A data engineering team is building a real-time data processing pipeline where application logs stored in S3 need to be filtered and compressed using a custom algorithm before being delivered to downstream applications. The team wants to maintain the standard S3 API interface while adding this transformation capability.\n\nWhich solution would best meet these requirements?",
            "options": [
                "Use Amazon Athena with custom SQL queries",
                "Implement client-side filtering with the S3 GET API",
                "Configure S3 Object Lambda with a custom transformation function",
                "Use Amazon EMR with Apache Spark for data transformation",
            ],
            "correct": [2],
            "explanations": [
                ["Athena is for ad-hoc analytics, not real-time transformation", "Doesn't maintain standard S3 API interface", "Not designed for per-request data transformation", "Adds query overhead for simple filtering"],
                ["Requires all clients to implement filtering logic", "Transfers full objects before filtering — wasteful", "Doesn't maintain transparent S3 API interface", "Pushes complexity to every downstream application"],
                ["Transforms data on-the-fly during S3 GET requests", "Maintains standard S3 API — transparent to clients", "Custom Lambda function applies filtering + compression", "Purpose-built for this exact use case"],
                ["EMR is for batch processing, not per-request transforms", "Requires cluster management — high overhead", "Doesn't maintain S3 API interface", "Over-engineered for filtering and compression"],
            ],
        },
        {
            "id": "q3",
            "title": "PII Detection & Masking",
            "question": "An e-commerce company collects data in an S3 data lake that includes PII. They need an automated solution to identify PII in new and existing data, provide an overview for the security audit team, and invoke a masking application in real time when PII is found.\n\nWhich solution meets these requirements with the LEAST operational overhead?",
            "options": [
                "Create an AWS Lambda function to analyze data for PII. Configure S3 bucket notifications to invoke Lambda when a new object is uploaded.",
                "Configure S3 bucket notifications with an EventBridge rule for new object uploads. Set the masking application as the target.",
                "Enable Amazon Macie. Create an EventBridge rule for Macie findings. Set the masking application as the target.",
                "Enable Amazon Macie. Create a Lambda function on a schedule to poll Macie findings and invoke the masking application.",
            ],
            "correct": [2],
            "explanations": [
                ["Requires building custom PII detection logic", "Only scans new uploads, not existing data", "No built-in reporting for security audit team", "High operational overhead to maintain"],
                ["No PII identification step at all", "Triggers masking on every upload blindly", "Wasteful — processes files without PII", "Doesn't meet the 'identify PII' requirement"],
                ["Macie is purpose-built for PII discovery in S3", "Scans both new and existing data", "EventBridge provides real-time event-driven response", "Built-in dashboard for security audit overview"],
                ["Polling adds latency — not real-time", "Scheduled Lambda adds operational overhead", "May miss findings between poll intervals", "Event-driven approach is simpler and faster"],
            ],
        },
        {
            "id": "q4",
            "title": "Redshift Data Sharing",
            "question": "Your company operates a large e-commerce platform with multiple Redshift data warehouses. The marketing team needs read access to specific data sets for personalized recommendations without compromising security. Data must be up-to-date and consistent across all regions.\n\nWhich approach would you recommend?",
            "options": [
                "Create a separate Redshift cluster and export data from existing clusters to the new cluster for the marketing team.",
                "Grant the marketing team's IAM roles access to existing Redshift clusters with specific SQL permissions.",
                "Have the marketing team spin up a new Redshift cluster and use Amazon Redshift data sharing to securely share data across clusters.",
                "Export data from Redshift to S3 using UNLOAD. Have the marketing team query S3 using Amazon Athena.",
            ],
            "correct": [2],
            "explanations": [
                ["Requires data export/import — not real-time", "Data becomes stale immediately after export", "Duplicate storage costs", "Doesn't maintain consistency across regions"],
                ["Direct access to production clusters", "Marketing queries may impact production workloads", "No workload isolation", "Violates security and governance requirements"],
                ["Data sharing provides live, read-only access", "No data copying — always up-to-date and consistent", "Workload isolation — separate compute for marketing", "Granular access control at schema/table level"],
                ["Data in S3 becomes stale after UNLOAD", "Requires ongoing sync pipeline to maintain freshness", "Athena query performance differs from Redshift", "Additional operational overhead for pipeline management"],
            ],
        },
        {
            "id": "q5",
            "title": "Credential Management",
            "question": "A company runs a cloud-based application in EC2 backed by RDS for Microsoft SQL Server. The application stores confidential information. The company wants to eliminate the risk of credential exposure.\n\nWhich solution will meet this requirement?",
            "options": [
                "Use AWS IAM database authentication to configure authentication to the RDS for Microsoft SQL Server database.",
                "Use AWS Systems Manager Parameter Store to store credentials. Configure automatic rotation every 30 days.",
                "Use AWS Security Token Service (AWS STS) to configure authentication to the RDS for Microsoft SQL Server database.",
                "Use AWS Secrets Manager to store the credentials. Configure automatic rotation every 30 days.",
            ],
            "correct": [3],
            "explanations": [
                ["IAM database auth is NOT supported for SQL Server", "Only works with MySQL, PostgreSQL, and MariaDB on RDS", "Cannot be used for this specific database engine", "Common exam trap — know which engines support IAM auth"],
                ["Parameter Store can store credentials securely", "Does NOT have built-in automatic rotation", "Rotation requires custom Lambda implementation", "More operational overhead than Secrets Manager"],
                ["STS provides temporary security credentials for AWS services", "Cannot directly authenticate to RDS databases", "Designed for assuming IAM roles, not DB authentication", "Wrong service for database credential management"],
                ["Secrets Manager is purpose-built for credential management", "Native automatic rotation for RDS databases", "Supports SQL Server credential rotation out of the box", "Eliminates credential exposure with zero custom code"],
            ],
        },
    ],
}

# ============================================================
# SESSION 5 QUESTIONS
# ============================================================
SESSIONS[5] = {
    "base_dir": "session5/exam-strategy",
    "questions": [
        {
            "id": "q1",
            "title": "Redshift Data Sharing (Cross-Account)",
            "question": "A data engineer needs to distribute processed data from a central Redshift ETL cluster to multiple BI teams across different AWS accounts and regions. Teams need different access levels, real-time updates without data copies, cost-effectiveness, and workload isolation.\n\nWhat is the most suitable approach?",
            "options": [
                "Set up cross-region snapshot copies and restore them in each region.",
                "Use Amazon Redshift data sharing with a central ETL cluster.",
                "Implement real-time replication using AWS DMS.",
                "Create separate ETL processes in each region.",
            ],
            "correct": [1],
            "explanations": [
                ["Snapshots create point-in-time copies — not real-time", "Requires restore in each region — operational overhead", "Data becomes stale between snapshot intervals", "Storage duplication across regions increases cost"],
                ["Live data sharing — no copies, always current", "Supports cross-account and cross-region sharing", "Granular access control per consumer cluster", "Workload isolation with separate compute per team"],
                ["DMS is for database migration, not data distribution", "Adds replication lag and operational complexity", "Not designed for multi-consumer read access patterns", "Over-engineered for this use case"],
                ["Duplicate ETL pipelines in every region", "Highest operational overhead and cost", "Data inconsistency risk across regions", "Violates the 'no data copies' requirement"],
            ],
        },
        {
            "id": "q2",
            "title": "Redshift Encryption Modification",
            "question": "A data engineer needs to enable encryption on an existing unencrypted Redshift cluster using RA3 nodes. The cluster processes critical financial data 24/7. They need to minimize downtime and maintain query capabilities during encryption.\n\nWhich statement accurately describes the capabilities for this encryption modification?",
            "options": [
                "The cluster must be completely stopped during encryption, and no queries can be run until the process is complete.",
                "Only read queries can be executed during the encryption process, and write operations must be suspended.",
                "Both read and write queries can continue during encryption with RA3 nodes, but elastic resize should be performed after encryption starts.",
                "Both read and write queries can continue during encryption with RA3 nodes, but elastic resize should be performed before encryption starts.",
            ],
            "correct": [3],
            "explanations": [
                ["RA3 nodes support in-place encryption without stopping", "Cluster remains fully operational during encryption", "No downtime required for RA3 node types", "This was true for older node types, not RA3"],
                ["RA3 supports both read AND write during encryption", "No need to suspend write operations", "Full query capability is maintained", "Overly restrictive — doesn't reflect RA3 capabilities"],
                ["Read and write continuing is correct for RA3", "But elastic resize AFTER encryption causes issues", "Resize during encryption can conflict with the process", "Resize should be done BEFORE to avoid complications"],
                ["RA3 nodes support full read/write during encryption", "Elastic resize before encryption avoids conflicts", "Encryption process runs in background on RA3", "Best practice: resize first, then enable encryption"],
            ],
        },
        {
            "id": "q3",
            "title": "Lake Formation Permissions",
            "question": "A retail company manages a data lake with sensitive customer data. They use complex IAM roles and S3 bucket policies for access. They're experiencing challenges with permission management, especially column-level access for different analyst groups.\n\nWhich solution would best address their needs following AWS best practices?",
            "options": [
                "Create additional IAM roles with specific S3 bucket policies for each analyst group.",
                "Implement AWS Lake Formation to manage permissions on Data Catalog resources and underlying S3 data.",
                "Use AWS Secrets Manager to store access credentials for different analyst groups.",
                "Configure VPC endpoints with custom route tables for controlled access to S3 buckets.",
            ],
            "correct": [1],
            "explanations": [
                ["More IAM roles adds to existing complexity", "S3 bucket policies cannot provide column-level access", "Doesn't simplify — makes the problem worse", "IAM + bucket policies is what they're trying to move away from"],
                ["Lake Formation provides centralized permission management", "Supports column-level and row-level access control", "Simplifies complex IAM/S3 policy combinations", "AWS recommended approach for data lake security"],
                ["Secrets Manager stores credentials, not access policies", "Doesn't provide column-level access control", "Wrong tool — designed for secrets, not authorization", "Doesn't address the permission management challenge"],
                ["VPC endpoints control network path, not data access", "Cannot provide column-level permissions", "Route tables don't filter data content", "Addresses network security, not data authorization"],
            ],
        },
        {
            "id": "q4",
            "title": "Glue Data Catalog Encryption",
            "question": "A financial services company needs to implement encryption for their AWS Glue Data Catalog. They want full control over encryption keys, ability to rotate cryptographic material, modify key policies, and ensure audit logging support.\n\nWhat should they implement?",
            "options": [
                "Customer managed keys with CloudTrail logging enabled.",
                "AWS managed keys with CloudTrail logging enabled.",
                "AWS managed keys with encryption context only.",
                "Custom encryption keys stored in AWS Secrets Manager.",
            ],
            "correct": [0],
            "explanations": [
                ["Customer managed keys (CMK) provide full control", "Can rotate cryptographic material on schedule", "Can modify key policies and grant access", "CloudTrail logs all key usage for audit compliance"],
                ["AWS managed keys don't allow key policy modification", "Cannot control rotation schedule", "Limited control over key lifecycle", "CloudTrail logging works but key control is insufficient"],
                ["AWS managed keys lack required control", "Encryption context alone doesn't provide audit logging", "Cannot rotate or modify key policies", "Doesn't meet the full control requirement"],
                ["Secrets Manager stores secrets, not encryption keys", "KMS is the correct service for encryption key management", "Cannot integrate with Glue Data Catalog encryption", "Wrong service for this use case entirely"],
            ],
        },
        {
            "id": "q5",
            "title": "S3 Dual-Layer Encryption",
            "question": "A healthcare company needs enhanced security for storing sensitive patient data in S3. Compliance requires multiple layers of encryption for all data at rest. The engineer wants to enforce dual-layer encryption for all objects uploaded to an S3 bucket.\n\nWhich combination of services and configurations would meet these requirements?",
            "options": [
                "Configure SSE-S3 encryption with S3 Bucket Keys enabled.",
                "Use client-side encryption with AWS KMS keys before uploading to S3.",
                "Implement DSSE-KMS encryption using AWS KMS keys in the same region as the bucket.",
                "Enable both SSE-KMS and SSE-S3 encryption on the bucket simultaneously.",
            ],
            "correct": [2],
            "explanations": [
                ["SSE-S3 is single-layer encryption only", "Bucket Keys optimize KMS costs but don't add layers", "Doesn't meet dual-layer encryption requirement", "Only one encryption layer applied"],
                ["Client-side + server-side could provide two layers", "But this is client-side only — not enforced at bucket level", "Cannot enforce via bucket policy alone", "Operational overhead managing client-side encryption"],
                ["DSSE-KMS is purpose-built dual-layer server-side encryption", "Two independent layers of encryption applied automatically", "KMS keys in same region required for DSSE-KMS", "Can be enforced via bucket policy — meets compliance"],
                ["S3 doesn't support enabling two SSE types simultaneously", "Must choose one server-side encryption method per object", "This configuration is not technically possible", "Invalid approach — S3 rejects conflicting encryption settings"],
            ],
        },
    ],
}


# ============================================================
# GENERATION LOGIC
# ============================================================
import subprocess
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Import the diagram generation and PPTX creation from session1's script
# We'll reuse the same PPTX creation logic

TIMER_PATH = "session1/exam-strategy/countdown_2min.gif"

NAVY = RGBColor(0x0d, 0x1b, 0x2a)
DARK_NAVY = RGBColor(0x09, 0x12, 0x1c)
ACCENT_BLUE = RGBColor(0x1e, 0x88, 0xe5)
CORRECT_GREEN = RGBColor(0x00, 0xc8, 0x53)
INCORRECT_RED = RGBColor(0xff, 0x17, 0x44)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf8, 0xf9, 0xfa)
MID_GRAY = RGBColor(0x90, 0xa4, 0xae)
DARK_TEXT = RGBColor(0x26, 0x32, 0x38)
TIMER_BG = RGBColor(0x1a, 0x23, 0x7e)
TIMER_RING = RGBColor(0x00, 0xe6, 0x76)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def create_question_slide(prs, title, question_text, options):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = NAVY
    fill.gradient_stops[1].color.rgb = DARK_NAVY

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.08), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_BLUE
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(6.5), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.font.name = "Segoe UI"

    q_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.8), Inches(7.2), Inches(1.5))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question_text
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    opt_top = 2.4
    for i, opt in enumerate(options):
        letter = chr(65 + i)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(opt_top - 0.02), Inches(0.32), Inches(0.32))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_BLUE
        badge.line.fill.background()
        badge_tf = badge.text_frame
        badge_tf.margin_top = Pt(0)
        badge_tf.margin_bottom = Pt(0)
        badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = badge_tf.paragraphs[0]
        bp.text = letter
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = WHITE
        bp.alignment = PP_ALIGN.CENTER

        opt_box = slide.shapes.add_textbox(Inches(0.95), Inches(opt_top), Inches(7.0), Inches(0.4))
        tf = opt_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = opt
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0xec, 0xef, 0xf1)
        p.font.name = "Segoe UI"
        opt_top += 0.58

    timer_card = add_rounded_rect(slide, Inches(7.7), Inches(0.3), Inches(2.1), Inches(2.3), TIMER_BG)
    timer_label = slide.shapes.add_textbox(Inches(7.8), Inches(0.4), Inches(1.9), Inches(0.35))
    tf = timer_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "\u23f1 TIME"
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = TIMER_RING
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    if os.path.exists(TIMER_PATH):
        slide.shapes.add_picture(TIMER_PATH, Inches(8.0), Inches(0.7), Inches(1.7), Inches(1.7))


def create_option_slide(prs, letter, option_text, diagram_path, explanation_bullets, is_correct):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = LIGHT_GRAY

    status_color = CORRECT_GREEN if is_correct else INCORRECT_RED
    status_text = "\u2713 CORRECT" if is_correct else "\u2717 INCORRECT"

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.9))
    banner.fill.solid()
    banner.fill.fore_color.rgb = NAVY
    banner.line.fill.background()

    badge = add_rounded_rect(slide, Inches(0.3), Inches(0.2), Inches(1.4), Inches(0.45), status_color)
    badge_tf = badge.text_frame
    badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = badge_tf.paragraphs[0]
    bp.text = status_text
    bp.font.size = Pt(10)
    bp.font.bold = True
    bp.font.color.rgb = WHITE
    bp.font.name = "Segoe UI"
    bp.alignment = PP_ALIGN.CENTER

    opt_box = slide.shapes.add_textbox(Inches(1.9), Inches(0.15), Inches(7.8), Inches(0.65))
    tf = opt_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = f"Option {letter}: {option_text}"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.name = "Segoe UI"

    if os.path.exists(diagram_path):
        with Image.open(diagram_path) as img:
            img_w, img_h = img.size
        max_w_in, max_h_in = 5.8, 3.4
        ratio = min(max_w_in / (img_w / 96.0), max_h_in / (img_h / 96.0))
        display_w = min((img_w / 96.0) * ratio, max_w_in)
        display_h = min((img_h / 96.0) * ratio, max_h_in)
        if display_w > max_w_in:
            display_h = (img_h / img_w) * max_w_in
            display_w = max_w_in
        if display_h > max_h_in:
            display_w = (img_w / img_h) * max_h_in
            display_h = max_h_in
        img_top = 1.1 + (3.6 - display_h) / 2
        slide.shapes.add_picture(diagram_path, Inches(0.2), Inches(img_top), Inches(display_w), Inches(display_h))

    panel = add_rounded_rect(slide, Inches(6.2), Inches(1.1), Inches(3.6), Inches(4.2), WHITE)
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(1.1), Inches(0.06), Inches(4.2))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = status_color
    accent_line.line.fill.background()

    why_box = slide.shapes.add_textbox(Inches(6.45), Inches(1.25), Inches(3.2), Inches(0.35))
    tf = why_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Why this is correct:" if is_correct else "Why this is incorrect:"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = status_color
    p.font.name = "Segoe UI"

    bullets_box = slide.shapes.add_textbox(Inches(6.45), Inches(1.65), Inches(3.2), Inches(3.5))
    tf = bullets_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(explanation_bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"\u2022 {bullet}"
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_TEXT
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)


def create_markdown(base_dir, q):
    """Create markdown file for a question."""
    q_dir = os.path.join(base_dir, q["id"])
    os.makedirs(q_dir, exist_ok=True)
    md_path = os.path.join(q_dir, f"{q['id']}.md")

    lines = [f"# {q['title']}\n\n{q['question']}\n"]
    for i, opt in enumerate(q["options"]):
        letter = chr(65 + i)
        lines.append(f"- {letter}) {opt}")
    lines.append("\n## Answer\n")
    correct_letters = [chr(65 + c) for c in q["correct"]]
    lines.append(f"**{'** and **'.join(correct_letters)}**\n")
    for c in q["correct"]:
        letter = chr(65 + c)
        lines.append(f"- **{letter}) {q['options'][c]}**")
        for bullet in q["explanations"][c]:
            lines.append(f"  - {bullet}")
    lines.append("\n## Why the other options are incorrect\n")
    for i, (opt, expl) in enumerate(zip(q["options"], q["explanations"])):
        if i not in q["correct"]:
            letter = chr(65 + i)
            lines.append(f"- **{letter}) {opt}**")
            for bullet in expl:
                lines.append(f"  - {bullet}")

    with open(md_path, "w") as f:
        f.write("\n".join(lines) + "\n")
    return md_path


def create_diagram_script(base_dir, q):
    """Create and run a diagram generation script for a question."""
    q_dir = os.path.join(base_dir, q["id"])
    os.makedirs(q_dir, exist_ok=True)

    # Build a simple diagram script for each option
    script_lines = [
        '"""Auto-generated diagram script."""',
        'import os, sys',
        'os.chdir(os.path.dirname(os.path.abspath(__file__)))',
        'from diagrams import Diagram, Cluster, Edge',
        'from diagrams.aws.analytics import KinesisDataStreams, Redshift, Glue, Athena',
        'from diagrams.aws.compute import Lambda, EC2',
        'from diagrams.aws.storage import S3, ElasticFileSystemEFS, ElasticBlockStoreEBS',
        'from diagrams.aws.database import Dynamodb, RDS',
        'from diagrams.aws.integration import Eventbridge',
        'from diagrams.aws.network import VPC, NATGateway',
        'from diagrams.aws.security import Shield',
        'from diagrams.aws.general import Users',
        '',
        'CORRECT_CLUSTER = {"bgcolor": "#e8f5e9", "style": "rounded", "pencolor": "#2e7d32", "penwidth": "2"}',
        'INCORRECT_CLUSTER = {"bgcolor": "#fbe9e7", "style": "rounded", "pencolor": "#c62828", "penwidth": "2"}',
        'NEUTRAL_CLUSTER = {"bgcolor": "#e3f2fd", "style": "rounded", "pencolor": "#1565c0", "penwidth": "1.5"}',
        '',
        'base_graph_attr = {"fontsize": "14", "bgcolor": "#fafafa", "pad": "0.8", "splines": "spline", "nodesep": "1.0", "ranksep": "1.0", "fontname": "Helvetica"}',
        '',
    ]

    # Generate a simple but meaningful diagram for each option
    for i, opt in enumerate(q["options"]):
        letter = chr(97 + i)  # a, b, c, d, e
        is_correct = i in q["correct"]
        cluster_style = "CORRECT_CLUSTER" if is_correct else "INCORRECT_CLUSTER"
        filename = f"option_{letter}"

        script_lines.append(f'with Diagram("", filename="{filename}", show=False, direction="LR", graph_attr=base_graph_attr):')
        script_lines.append(f'    with Cluster("{chr(65+i)}: {"Correct" if is_correct else "Incorrect"}", graph_attr={cluster_style}):')
        script_lines.append(f'        svc = S3("{chr(65+i)}")')
        script_lines.append('')

    script_lines.append(f'print("{q["id"]} diagrams generated.")')

    script_path = os.path.join(q_dir, "create_diagrams.py")
    with open(script_path, "w") as f:
        f.write("\n".join(script_lines) + "\n")

    # Run the script
    abs_script = os.path.abspath(script_path)
    abs_dir = os.path.abspath(q_dir)
    result = subprocess.run(["python3", abs_script], capture_output=True, text=True, cwd=abs_dir)
    if result.returncode != 0:
        print(f"  WARNING: Diagram generation failed for {q['id']}: {result.stderr[:200]}")
    return q_dir


def create_pptx(base_dir, q):
    """Create PowerPoint presentation for a question."""
    q_dir = os.path.join(base_dir, q["id"])
    os.makedirs(q_dir, exist_ok=True)

    prs = PptxPresentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    create_question_slide(prs, q["title"], q["question"], q["options"])

    for i, (opt, expl) in enumerate(zip(q["options"], q["explanations"])):
        letter = chr(65 + i)
        diagram_path = os.path.join(q_dir, f"option_{chr(97+i)}.png")
        is_correct = i in q["correct"]
        create_option_slide(prs, letter, opt, diagram_path, expl, is_correct)

    output_path = os.path.join(q_dir, f"{q['id']}_exam_strategy.pptx")
    prs.save(output_path)
    return output_path


# ============================================================
# MAIN EXECUTION
# ============================================================
if __name__ == "__main__":
    os.chdir(os.path.dirname(os.path.abspath(__file__)))

    for session_num, session_data in SESSIONS.items():
        base_dir = session_data["base_dir"]
        os.makedirs(base_dir, exist_ok=True)
        print(f"\n{'='*60}")
        print(f"SESSION {session_num}")
        print(f"{'='*60}")

        for q in session_data["questions"]:
            print(f"\n  Processing {q['id']}: {q['title']}")

            # 1. Create markdown
            md_path = create_markdown(base_dir, q)
            print(f"    ✓ Markdown: {md_path}")

            # 2. Create PowerPoint (diagrams already generated by generate_diagrams_s2_s5.py)
            pptx_path = create_pptx(base_dir, q)
            print(f"    ✓ PowerPoint: {pptx_path}")

    print(f"\n{'='*60}")
    print("ALL SESSIONS COMPLETE!")
    print(f"{'='*60}")
